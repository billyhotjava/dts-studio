#!/usr/bin/env python3
"""Generate DTS v3.0 Business Plan PPT — V1.0.0 edition."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# -- Colors --
BG = RGBColor(0x1A, 0x1A, 0x2E)
CARD = RGBColor(0x22, 0x22, 0x3A)
ACCENT = RGBColor(0x00, 0x96, 0xD6)
ACCENT2 = RGBColor(0x00, 0xC9, 0xA7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xAA, 0xAA, 0xBB)
LIGHT = RGBColor(0xE0, 0xE0, 0xE8)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
RED = RGBColor(0xE8, 0x4D, 0x4D)
GREEN = RGBColor(0x28, 0xC7, 0x6F)
PURPLE = RGBColor(0xBB, 0x86, 0xFC)
DARK_BLUE = RGBColor(0x0A, 0x3D, 0x5C)
DARK_RED = RGBColor(0x3A, 0x1A, 0x1A)
DARK_GREEN = RGBColor(0x0A, 0x3D, 0x2A)
DARK_PURPLE = RGBColor(0x2A, 0x1A, 0x3A)
YELLOW = RGBColor(0xFF, 0xD7, 0x00)


def bg(slide):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = BG

def txt(slide, l, t, w, h, text, sz=18, clr=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = clr; p.font.bold = bold
    p.font.name = "Microsoft YaHei"; p.alignment = align
    return tb

def card(slide, l, t, w, h, clr=CARD):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = clr; s.line.fill.background(); s.shadow.inherit = False
    return s

def bullets(slide, l, t, w, h, items, sz=16, clr=LIGHT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(sz); p.font.color.rgb = clr
        p.font.name = "Microsoft YaHei"; p.space_after = Pt(8)
    return tb

def bar(slide, l, t, w, h, clr):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = clr; s.line.fill.background()
    return s

def circle(slide, l, t, d, clr):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(d), Inches(d))
    s.fill.solid(); s.fill.fore_color.rgb = clr; s.line.fill.background()
    return s


# ================================================================
# Slide 1: Cover
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
bar(s, 1, 2.3, 1.5, 0.06, ACCENT)
txt(s, 1, 2.6, 11, 1.2, "Decision Twins System", 48, WHITE, True)
txt(s, 1, 3.7, 11, 0.8, "AI Decision Operating System -- V1.0.0", 28, ACCENT)
txt(s, 1, 4.6, 11, 0.6, "Ontology + Intent + Agent = Decision Twins", 20, GRAY)
txt(s, 1, 5.4, 8, 0.5, "2026-03  |  CONFIDENTIAL", 16, GRAY)


# ================================================================
# Slide 2: Executive Summary
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.8, 0.4, 12, 0.7, "Executive Summary", 36, WHITE, True)

card(s, 0.8, 1.3, 11.7, 2.0, DARK_BLUE)
txt(s, 1.3, 1.5, 11, 0.6,
    "把专家脑中的隐性经验，变成系统里能自动运转的「能力包」", 24, WHITE, True)
txt(s, 1.3, 2.1, 11, 1.0,
    "企业不缺数据，缺的是基于数据的决策能力。DTS 将专家 20 年经验通过 AI 对话提取、结构化编码、系统化运行，\n"
    "让用户用一句话完成「理解需求 -> 查数据 -> 分析 -> 图表 -> 报告 -> 行动建议」端到端决策闭环。", 17, LIGHT)

# Key metrics
metrics = [
    ("Y1 收入", "800 万", ACCENT),
    ("Y3 收入", "1 亿", GREEN),
    ("AI 技能", "57 个", ACCENT2),
    ("微服务", "25 个", ORANGE),
    ("融资需求", "2,000 万", PURPLE),
]
for i, (label, value, clr) in enumerate(metrics):
    x = 0.8 + i * 2.4
    card(s, x, 3.7, 2.1, 1.3)
    bar(s, x, 3.7, 2.1, 0.06, clr)
    txt(s, x + 0.2, 3.9, 1.7, 0.4, label, 15, GRAY)
    txt(s, x + 0.2, 4.3, 1.7, 0.5, value, 24, clr, True)

# Differentiation
txt(s, 0.8, 5.3, 12, 0.4, "核心差异化", 20, WHITE, True)
diffs = [
    ("vs 数据中台", "解决「有没有数据」\nDTS 解决「数据能不能做决策」", ACCENT),
    ("vs BI 工具", "展示数据让人看\nDTS 理解意图替人想、帮人做", ACCENT2),
    ("vs AI Agent", "通用 Agent 无行业深度\nDTS = Ontology + KnowHow", GREEN),
    ("vs Palantir", "重、贵、封闭\nDTS 轻量私有化、开放生态", ORANGE),
]
for i, (title, desc, clr) in enumerate(diffs):
    x = 0.8 + i * 3.1
    card(s, x, 5.8, 2.8, 1.5)
    txt(s, x + 0.2, 5.85, 2.4, 0.35, title, 15, clr, True)
    txt(s, x + 0.2, 6.25, 2.4, 0.9, desc, 14, LIGHT)


# ================================================================
# Slide 3: Know-how Value Proposition
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.8, 0.4, 12, 0.7, "Know-how 为什么值钱", 36, WHITE, True)
txt(s, 0.8, 1.1, 12, 0.5,
    "企业的核心竞争力，大多数存在几个人的脑子里", 18, GRAY)

# Left: without
card(s, 0.8, 1.8, 5.8, 2.5, DARK_RED)
txt(s, 1.1, 1.95, 5.2, 0.4, "没有 Know-how 的现实", 20, RED, True)
bullets(s, 1.1, 2.4, 5.2, 1.8, [
    "500 个变电站，只有 3 个真正有经验的调度专家",
    "资深工程师 8-10 年培养周期，退休即断层",
    "年均流动率 15-20%，知识归零、重复交学费",
    "决策质量取决于「今天谁值班」而非组织能力",
], 15, LIGHT)

# Right: with
card(s, 7.0, 1.8, 5.8, 2.5, DARK_GREEN)
txt(s, 7.3, 1.95, 5.2, 0.4, "Know-how 变成能力包之后", 20, GREEN, True)
bullets(s, 7.3, 2.4, 5.2, 1.8, [
    "3 个专家的判断力 -> 复制给 500 个站",
    "新人上手就有老专家水平的辅助决策",
    "专家退休，经验还在系统里继续运转",
    "每个项目积累的 Know-how 自动沉淀，越用越值钱",
], 15, LIGHT)

# Three values
txt(s, 0.8, 4.6, 12, 0.5, "三层价值", 20, ACCENT, True)
layers = [
    ("省人", "装上能力包\n这件事不再需要人盯", "直接减少\n重复劳动岗位", GREEN),
    ("提质", "每个新人都能像\n老专家一样判断", "错误率下降\n决策质量统一", ACCENT),
    ("保全", "专家可以退休\n能力留在系统里", "核心竞争力\n不随人员流失", ORANGE),
]
for i, (title, desc, val, clr) in enumerate(layers):
    x = 0.8 + i * 4.1
    card(s, x, 5.2, 3.8, 2.1)
    bar(s, x, 5.2, 3.8, 0.08, clr)
    txt(s, x + 0.3, 5.45, 3.2, 0.5, title, 22, clr, True)
    txt(s, x + 0.3, 5.95, 3.2, 0.8, desc, 15, LIGHT)
    txt(s, x + 0.3, 6.7, 3.2, 0.5, val, 13, GRAY)


# ================================================================
# Slide 4: Market Timing
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.8, 0.4, 12, 0.7, "为什么是现在", 36, WHITE, True)
txt(s, 0.8, 1.1, 12, 0.5, "三个条件在 2025-2026 年首次同时成熟", 18, GRAY)

conditions = [
    ("AI Agent 能力就绪", "大模型推理/规划/工具调用\n突破临界点\nAgent 从玩具变生产力",
     "Manus 8 月达 $100M ARR\n验证 Agent 市场需求", ACCENT),
    ("企业数据就绪", "过去 5-10 年投入数百亿\n建数据湖/数仓/ETL/BI\nROI 普遍不达预期",
     "不是没数据，是有数据没价值\nDTS = 最后一公里", ACCENT2),
    ("客户意愿就绪", "决策者亲身体验 AI 能力\n首次主动寻求\nAI+业务决策的落地方案",
     "从被动接受 IT 提案\n转为主动推动", GREEN),
]
for i, (title, desc, note, clr) in enumerate(conditions):
    x = 0.8 + i * 4.1
    card(s, x, 1.8, 3.8, 3.5)
    bar(s, x, 1.8, 3.8, 0.08, clr)
    txt(s, x + 0.3, 2.05, 3.2, 0.5, title, 22, clr, True)
    txt(s, x + 0.3, 2.6, 3.2, 1.2, desc, 16, LIGHT)
    txt(s, x + 0.3, 3.8, 3.2, 1.0, note, 13, GRAY)

card(s, 0.8, 5.7, 11.7, 1.3, DARK_BLUE)
txt(s, 1.2, 5.9, 11, 0.5, "窗口期判断", 22, WHITE, True)
txt(s, 1.2, 6.4, 11, 0.5,
    "太早（2023）：AI 不可靠，客户不信  |  "
    "太晚（2028）：大厂做了，赛道拥挤  |  "
    "现在（2026）：客户有数据、有预算、有意愿，但没有好产品", 15, LIGHT)


# ================================================================
# Slide 5: Three Pain Points
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.8, 0.4, 12, 0.7, "企业决策的三大痛点", 36, WHITE, True)

pains = [
    ("专家依赖症", "核心决策高度依赖\n少数资深专家\n电网调度员 8-10 年培养\n决策依据在「脑子里」",
     "波兰尼：We know more\nthan we can tell", RED),
    ("经验流失危机", "核心技术岗年均\n流动率 15-20%\n每流失一位资深工程师\n需 2-3 年重新培养",
     "知识归零\n重复交学费", ORANGE),
    ("决策质量不一致", "老手与新手分析\n质量差距达 10 倍\n决策取决于「今天谁值班」\n而非组织能力",
     "数千万数据基础设施\n产出仍靠个人能力", YELLOW),
]
for i, (title, desc, note, clr) in enumerate(pains):
    x = 0.8 + i * 4.1
    card(s, x, 1.3, 3.8, 4.0)
    bar(s, x, 1.3, 3.8, 0.08, clr)
    circle(s, x + 0.3, 1.6, 0.6, clr)
    txt(s, x + 0.3, 1.65, 0.6, 0.5, str(i+1), 24, WHITE, True, PP_ALIGN.CENTER)
    txt(s, x + 1.1, 1.65, 2.5, 0.5, title, 22, WHITE, True)
    txt(s, x + 0.3, 2.3, 3.2, 1.6, desc, 16, LIGHT)
    txt(s, x + 0.3, 4.0, 3.2, 0.9, note, 14, GRAY)

card(s, 0.8, 5.7, 11.7, 1.5, DARK_BLUE)
txt(s, 1.2, 5.9, 11, 0.5, "市场上缺什么？", 22, ACCENT, True)
txt(s, 1.2, 6.4, 11, 0.7,
    "缺一个同时具备的产品：连接企业数据 + 沉淀行业 Know-how + AI 多步决策执行 + 安全合规 + 可插拔定制\n"
    "每一项单独都有人做，但组合在一起的没有 -- 这就是 DTS 的位置", 17, LIGHT)


# ================================================================
# Slide 6: TAM / SAM / SOM
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.8, 0.4, 12, 0.7, "市场规模 TAM / SAM / SOM", 36, WHITE, True)

# TAM
card(s, 0.8, 1.3, 3.6, 2.8)
bar(s, 0.8, 1.3, 3.6, 0.08, ACCENT)
txt(s, 1.1, 1.55, 3.0, 0.4, "TAM", 24, ACCENT, True)
txt(s, 1.1, 2.0, 3.0, 0.4, "2,500 亿", 28, WHITE, True)
txt(s, 1.1, 2.5, 3.0, 1.2,
    "中国企业 AI 应用市场\n2026 年 年增速 35%+\n决策智能细分 ~800 亿", 15, LIGHT)

# SAM
card(s, 4.7, 1.3, 3.6, 2.8)
bar(s, 4.7, 1.3, 3.6, 0.08, ACCENT2)
txt(s, 5.0, 1.55, 3.0, 0.4, "SAM", 24, ACCENT2, True)
txt(s, 5.0, 2.0, 3.0, 0.4, "310 亿", 28, WHITE, True)
txt(s, 5.0, 2.5, 3.0, 1.2,
    "能源 ~120 亿\n科研院所 ~40 亿\n先进制造 ~150 亿", 15, LIGHT)

# SOM
card(s, 8.6, 1.3, 3.6, 2.8)
bar(s, 8.6, 1.3, 3.6, 0.08, GREEN)
txt(s, 8.9, 1.55, 3.0, 0.4, "SOM (3 年)", 24, GREEN, True)
txt(s, 8.9, 2.0, 3.0, 0.4, "1 亿", 28, WHITE, True)
txt(s, 8.9, 2.5, 3.0, 1.2,
    "Y1: 800 万 (2-3 客户)\nY2: 3,500 万 (8-12 客户)\nY3: 1 亿 (20+ 客户)", 15, LIGHT)

# Target industries
txt(s, 0.8, 4.4, 12, 0.5, "目标行业", 22, WHITE, True)
industries = [
    ("能源/电网", "200-500 万/年\n调度优化 | 设备维保\n负荷预测 | 安全监控", ACCENT),
    ("科研院所", "50-150 万/年\n实验数据管理\n智能分析 | 论文辅助", ACCENT2),
    ("先进制造", "100-300 万/年\n质量决策 | 排产优化\n供应链管理", GREEN),
    ("新能源", "100-200 万/年\n风电/光伏/储能\n运维决策", ORANGE),
]
for i, (name, desc, clr) in enumerate(industries):
    x = 0.8 + i * 3.1
    card(s, x, 5.0, 2.8, 2.2)
    bar(s, x, 5.0, 2.8, 0.08, clr)
    txt(s, x + 0.2, 5.2, 2.4, 0.4, name, 18, clr, True)
    txt(s, x + 0.2, 5.65, 2.4, 1.2, desc, 14, LIGHT)


# ================================================================
# Slide 7: Decision Twin Concept
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.8, 0.4, 12, 0.7, "Decision Twin -- 决策过程的数字孪生", 36, WHITE, True)
txt(s, 0.8, 1.1, 12, 0.5,
    "不是模拟物理设备，而是复制专家的决策能力", 18, GRAY)

# Comparison table
headers = ["维度", "Digital Twin (物理孪生)", "Decision Twin (决策孪生)"]
rows_dt = [
    ("孪生对象", "物理设备的结构和状态", "专家的决策过程和判断逻辑"),
    ("核心价值", "模拟物理行为", "复制决策能力"),
    ("关键输入", "传感器数据、CAD 模型", "隐性知识、业务规则、行业经验"),
    ("关键输出", "设备状态预测", "决策建议、行动方案、分析报告"),
    ("技术基础", "IoT + 仿真引擎", "Ontology + AI Agent + KnowHow"),
]
cw = [2.0, 4.5, 4.5]
# Headers
x_pos = 0.8
for j, h in enumerate(headers):
    txt(s, x_pos, 1.7, cw[j], 0.4, h, 15, ACCENT, True)
    x_pos += cw[j]
# Rows
for ri, (dim, dt, dct) in enumerate(rows_dt):
    y = 2.15 + ri * 0.55
    if ri % 2 == 0:
        bar(s, 0.8, y, sum(cw), 0.5, RGBColor(0x1E, 0x1E, 0x32))
    x_pos = 0.8
    for j, v in enumerate([dim, dt, dct]):
        c = GRAY if j == 0 else (RED if j == 1 else GREEN)
        txt(s, x_pos, y + 0.05, cw[j], 0.4, v, 14, c)
        x_pos += cw[j]

card(s, 0.8, 5.1, 11.7, 2.0, DARK_BLUE)
txt(s, 1.2, 5.3, 11, 0.5,
    "一句话定义", 22, WHITE, True)
txt(s, 1.2, 5.8, 11, 1.0,
    "Decision Twin 把企业的数据资产、业务规则、行业知识编织成一个能替决策者「想」的系统。\n"
    "不是替代人决策，而是让每个人都拥有专家级的决策能力。\n\n"
    "核心公式：Ontology (认知基座) + Intent (意图理解) + Agent (自主执行)", 17, LIGHT)


# ================================================================
# Slide 8: vs Traditional Solutions
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.8, 0.4, 12, 0.7, "为什么现有方案解决不了", 36, WHITE, True)

alternatives = [
    ("通用 AI\nChatGPT/文心", "没有企业数据\n不能离线部署\n不沉淀行业知识\n数据安全不可控"),
    ("BI 工具\n帆软/Tableau", "只能「看」不能「做」\n没有 AI 决策能力\n没有行业 Know-how"),
    ("AI 搭建平台\nCoze/Dify", "轻量 Agent 搭建\n没有数据治理\n无行业深度\n不能离线部署"),
    ("大厂平台\n阿里/华为", "做通用平台\n不深耕行业 Know-how\n政企客户不想被锁定"),
    ("外包定制\nSI 做项目", "做一个项目的知识\n不能复用到下一个\n乙方走了系统就废了"),
]
for i, (name, problem) in enumerate(alternatives):
    x = 0.8 + i * 2.46
    card(s, x, 1.3, 2.2, 3.2)
    txt(s, x + 0.15, 1.4, 1.9, 0.7, name, 15, WHITE, True)
    txt(s, x + 0.15, 2.2, 1.9, 1.8, problem, 14, RED)

card(s, 0.8, 4.9, 11.7, 2.2, DARK_BLUE)
txt(s, 1.2, 5.1, 11, 0.5, "DTS 的位置", 22, ACCENT, True)
txt(s, 1.2, 5.6, 11, 1.2,
    "同时具备五个要素：连接企业数据 + 沉淀行业 Know-how + AI 多步决策执行 + 安全合规 + 可插拔定制\n\n"
    "不和 ChatGPT 比通用能力 | 不和帆软比报表数量 | 做「AI + 行业数据 + Know-how 沉淀」的交叉地带", 17, LIGHT)


# ================================================================
# Slide 9: Product Architecture (Three Pillars)
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.8, 0.4, 12, 0.7, "产品架构：三支柱", 36, WHITE, True)
txt(s, 0.8, 1.1, 12, 0.5,
    "设计、运行、应用三个关注点清晰分离", 18, GRAY)

pillars = [
    ("dts-studio", "设计中心", "Ontology 建模\nSkill 开发\nAppPack 打包发布\n统一 IDE", ACCENT),
    ("dts-stack", "Agent 能力中心", "25 个微服务\n57 个 AI 技能\nDAP 协议运行时\n五条铁律保障", ACCENT2),
    ("app-stack", "客户 Agent 中心", "行业 AppPack\n客户定制 Pack\nKnowHow 资产\nISV 开发的 Pack", GREEN),
]
for i, (name, subtitle, desc, clr) in enumerate(pillars):
    x = 0.8 + i * 4.1
    card(s, x, 1.7, 3.8, 3.0)
    bar(s, x, 1.7, 3.8, 0.08, clr)
    txt(s, x + 0.3, 1.95, 3.2, 0.4, name, 22, clr, True)
    txt(s, x + 0.3, 2.4, 3.2, 0.4, subtitle, 16, WHITE)
    txt(s, x + 0.3, 2.85, 3.2, 1.6, desc, 15, LIGHT)

# DAP connection
card(s, 0.8, 5.0, 11.7, 0.7, DARK_PURPLE)
txt(s, 1.2, 5.1, 11, 0.5,
    "DAP 协议贯穿三支柱 -- 定义 AI Agent 如何理解物理世界并执行决策", 18, PURPLE, True, PP_ALIGN.CENTER)

# Three-language collaboration
txt(s, 0.8, 5.9, 12, 0.4, "三语言协作", 20, WHITE, True)
langs = [
    ("Python", "10 服务", "AI 核心层\nAgent/Intent/RAG/Eval", ACCENT),
    ("Java", "10 服务", "业务平台层\nCRUD/RBAC/安全/存储", ACCENT2),
    ("Go", "2 服务", "基础设施层\nK8s Operator/CLI", GREEN),
    ("React", "3 应用", "前端层\nAdmin/Cortex/Pilot", ORANGE),
]
for i, (lang, count, desc, clr) in enumerate(langs):
    x = 0.8 + i * 3.1
    card(s, x, 6.3, 2.8, 1.1)
    txt(s, x + 0.2, 6.35, 1.2, 0.3, lang, 16, clr, True)
    txt(s, x + 1.5, 6.35, 1.1, 0.3, count, 14, GRAY)
    txt(s, x + 0.2, 6.7, 2.4, 0.6, desc, 12, LIGHT)


# ================================================================
# Slide 10: DAP Protocol
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.8, 0.4, 12, 0.7, "DAP 协议 -- 产品核心护城河", 36, WHITE, True)
txt(s, 0.8, 1.1, 12, 0.5,
    "DTS Agent Protocol: 定义 AI Agent 如何理解物理世界并执行决策", 18, GRAY)

dap_layers = [
    ("L5", "Audit Protocol", "决策可追溯性保障\nCloudEvents 标准, append-only", PURPLE),
    ("L4", "Security Protocol", "身份与权限在 Agent 间传递\nJWT + 零信任校验", RED),
    ("L3", "Intent Protocol", "自然语言 -> 结构化意图 -> 执行计划\n支持条件执行和风险分级", ORANGE),
    ("L2", "Skill Protocol", "AI 技能定义/注册/调用/编排\ngRPC 标准接口 + 流式响应", ACCENT2),
    ("L1", "Ontology Protocol", "物理世界 -> 数字映射标准格式\nObjectType/Relationship/Metric", ACCENT),
    ("L0", "Transport", "gRPC (控制面) + Kafka (数据面) + SSE (流式面)", GRAY),
]
for i, (layer, name, desc, clr) in enumerate(dap_layers):
    y = 1.7 + i * 0.82
    card(s, 0.8, y, 11.7, 0.72)
    bar(s, 0.8, y, 0.12, 0.72, clr)
    txt(s, 1.2, y + 0.08, 0.8, 0.3, layer, 16, clr, True)
    txt(s, 2.0, y + 0.08, 2.8, 0.3, name, 16, WHITE, True)
    txt(s, 5.0, y + 0.08, 7.2, 0.55, desc, 14, LIGHT)

# Why moat
card(s, 0.8, 6.7, 11.7, 0.7, DARK_BLUE)
txt(s, 1.2, 6.75, 11, 0.55,
    "代码可以被重写，协议定义的生态不能 | 客户迁移成本 = 丢失核心竞争力 | "
    "类比：Android 护城河不是 AOSP 代码，而是 Google Play 生态协议", 14, LIGHT)


# ================================================================
# Slide 11: Knowledge Strategy (Polanyi + LeCun)
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.8, 0.4, 12, 0.7, "知识战略：Polanyi + LeCun", 36, WHITE, True)
txt(s, 0.8, 1.1, 12, 0.5,
    "传统数据平台假设知识 = 可量化数据，但真正驱动决策的知识大部分不是数据", 16, GRAY)

# Two sources
card(s, 0.8, 1.7, 5.5, 1.5)
bar(s, 0.8, 1.7, 0.12, 1.5, ACCENT)
txt(s, 1.2, 1.8, 5.0, 0.4, "Michael Polanyi (波兰尼)", 18, ACCENT, True)
txt(s, 1.2, 2.2, 5.0, 0.8,
    "缄默知识理论: We know more than we can tell\n企业最有价值的知识在专家脑子里，传统系统无法存储", 14, LIGHT)

card(s, 6.7, 1.7, 5.5, 1.5)
bar(s, 6.7, 1.7, 0.12, 1.5, ACCENT2)
txt(s, 7.1, 1.8, 5.0, 0.4, "Yann LeCun (2026)", 18, ACCENT2, True)
txt(s, 7.1, 2.2, 5.0, 0.8,
    "Beyond Language Modeling: 语言只是物理世界的投影\n通用预训练 + 1% 领域数据即可涌现理解能力", 14, LIGHT)

# Five-layer pyramid
txt(s, 0.8, 3.5, 12, 0.4, "五层知识金字塔", 20, WHITE, True)
pyramid = [
    ("L4", "世界模型", "Agent 预测操作后果: State+Action->Next State", PURPLE, "远期"),
    ("L3", "隐性模式", "专家没说但一直在用的规则，多模态行为推断", ORANGE, "中期"),
    ("L2", "潜在知识", "专家能说但没人问过的判断逻辑 (Gen-Ba 碎片提取)", ACCENT2, "短期"),
    ("L1", "显性结构", "数据字典、行业标准、系统表结构 (传统 Ontology)", ACCENT, "当前"),
    ("L0", "多模态感知", "视觉+声音+时序+文本 (身体化知识)", GRAY, "基础"),
]
for i, (layer, name, desc, clr, timeline) in enumerate(pyramid):
    y = 3.95 + i * 0.6
    w = 11.7 - i * 0.6
    x = 0.8 + i * 0.3
    card(s, x, y, w, 0.52)
    bar(s, x, y, 0.1, 0.52, clr)
    txt(s, x + 0.3, y + 0.05, 0.5, 0.3, layer, 13, clr, True)
    txt(s, x + 0.9, y + 0.05, 1.8, 0.3, name, 14, WHITE, True)
    txt(s, x + 2.8, y + 0.05, w - 4.5, 0.4, desc, 13, LIGHT)
    txt(s, x + w - 1.5, y + 0.05, 1.3, 0.3, timeline, 12, clr)

txt(s, 0.8, 7.0, 12, 0.4,
    "传统 BI/数据中台只覆盖 L1 | DTS 目标: 逐步覆盖 L0-L4，构建企业对物理世界的完整认知", 15, GRAY)


# ================================================================
# Slide 12: KnowHow Extraction Pipeline
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.8, 0.4, 12, 0.7, "KnowHow 提取流水线", 36, WHITE, True)
txt(s, 0.8, 1.1, 12, 0.5,
    "超越 Nonaka SECI: 在决策瞬间自然追问，而非系统性访谈", 18, GRAY)

phases = [
    ("Phase 1", "影子观察", "AI Agent 静默观察\n专家日常决策", "不干扰工作流", ACCENT),
    ("Phase 2", "碎片外化", "决策瞬间自然追问\n「为什么选择提前维保？」", "容许不完整\n容许矛盾", ACCENT2),
    ("Phase 3", "松散组合", "LLM 将碎片聚类\n去矛盾、补全\n生成 KnowHow 草稿", "追求「够用的近似」\n非完美转化", GREEN),
    ("Phase 4", "协作内化", "KnowHow 转为 Skill\n与专家并行运行\n持续校正", "评测套件量化准确率", ORANGE),
]
for i, (phase, name, desc, principle, clr) in enumerate(phases):
    x = 0.8 + i * 3.1
    card(s, x, 1.8, 2.8, 3.5)
    bar(s, x, 1.8, 2.8, 0.08, clr)
    txt(s, x + 0.2, 2.0, 2.4, 0.3, phase, 14, clr, True)
    txt(s, x + 0.2, 2.35, 2.4, 0.4, name, 20, WHITE, True)
    txt(s, x + 0.2, 2.85, 2.4, 1.2, desc, 15, LIGHT)
    txt(s, x + 0.2, 4.1, 2.4, 0.8, principle, 13, GRAY)
    if i < 3:
        txt(s, x + 2.6, 2.8, 0.5, 0.4, "->", 20, GRAY)

card(s, 0.8, 5.6, 11.7, 1.6, DARK_BLUE)
txt(s, 1.2, 5.75, 11, 0.5,
    "与 SECI 模型的本质区别", 20, WHITE, True)
txt(s, 1.2, 6.2, 11, 0.8,
    "SECI 假设隐性知识能被「完全转化」为显性知识 -- 这是认识论幻觉。\n"
    "DTS 承认隐性知识有不可消除的缄默维度，只追求「够用的近似」，\n"
    "剩余不确定性由人工介入 (HITL) 覆盖 -- 这正是铁律 #1 的认识论基础。", 16, LIGHT)


# ================================================================
# Slide 13: Business Model (Three Revenue Layers)
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.8, 0.4, 12, 0.7, "商业模式：知识资产平台", 36, WHITE, True)
txt(s, 0.8, 1.1, 12, 0.4, "三层收入 -- 不是 SaaS 卖功能，而是让知识流动产生价值", 18, GRAY)

rev_layers = [
    ("第一层: 平台基础费", "35%", "按部署规模年费\n标准版 30-80 万\n企业版 100-300 万\n不按人头 -> 鼓励全员使用", ACCENT),
    ("第二层: AppPack 市场", "30%", "行业能力包按需购买\n10-50 万/个/年\n客户/ISV/专家创建\n平台抽成 30%", ACCENT2),
    ("第三层: KnowHow 交易", "10%", "专家经验 -> 结构化 KnowHow\n发布到市场，按使用量分润\n知识资产的 App Store\n越用越值钱", GREEN),
]
for i, (name, pct, desc, clr) in enumerate(rev_layers):
    x = 0.8 + i * 4.1
    card(s, x, 1.6, 3.8, 3.5)
    bar(s, x, 1.6, 3.8, 0.08, clr)
    txt(s, x + 0.3, 1.85, 3.2, 0.4, name, 18, clr, True)
    txt(s, x + 0.3, 2.3, 3.2, 0.3, f"Year 3 占比 {pct}", 14, GRAY)
    txt(s, x + 0.3, 2.7, 3.2, 2.0, desc, 15, LIGHT)

# + Professional services
card(s, 0.8, 5.3, 5.5, 1.0)
bar(s, 0.8, 5.3, 0.12, 1.0, ORANGE)
txt(s, 1.2, 5.35, 5.0, 0.35, "专业服务 (Y3: 25%)", 17, ORANGE, True)
txt(s, 1.2, 5.75, 5.0, 0.4, "实施+定制+培训 | 行业深度咨询 | 5 天 POV 交付", 14, LIGHT)

# Flywheel
card(s, 7.0, 5.3, 5.5, 1.0, DARK_BLUE)
txt(s, 7.3, 5.35, 5.0, 0.35, "「越用越值钱」飞轮", 17, WHITE, True)
txt(s, 7.3, 5.75, 5.0, 0.4,
    "创建能力包 -> 知识沉淀 -> 不可替代 -> 续费率升 -> 吸引创作者", 14, LIGHT)

# Customer tiers
txt(s, 0.8, 6.5, 12, 0.4, "客户分层", 18, WHITE, True)
tiers = [
    ("旗舰客户", "央企/大型国企", "200-500 万/年", ACCENT),
    ("标准客户", "中大型企业", "50-150 万/年", ACCENT2),
    ("中小客户", "数据驱动型中小", "10-30 万/年", GREEN),
]
for i, (tier, desc, price, clr) in enumerate(tiers):
    x = 0.8 + i * 4.1
    card(s, x, 6.9, 3.8, 0.5)
    txt(s, x + 0.2, 6.93, 1.2, 0.3, tier, 14, clr, True)
    txt(s, x + 1.5, 6.93, 1.2, 0.3, desc, 13, GRAY)
    txt(s, x + 2.8, 6.93, 1.0, 0.3, price, 13, WHITE, True)


# ================================================================
# Slide 14: Go-to-Market Strategy
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.8, 0.4, 12, 0.7, "落地策略", 36, WHITE, True)

# Four phases
txt(s, 0.8, 1.15, 12, 0.4, "四阶段路线图", 20, WHITE, True)
roadmap = [
    ("Phase 1", "影子先行", "0-6 月", "显性 Ontology\n+ Intent + Skill\n全链路跑通", "L1", ACCENT),
    ("Phase 2", "碎片积累", "6-12 月", "KnowHow 提取上线\n2-3 个标杆客户", "L0-L2", ACCENT2),
    ("Phase 3", "模式涌现", "12-18 月", "隐性模式推断\nAppPack 市场 Beta\n5-8 个客户", "L3", GREEN),
    ("Phase 4", "寓居深化", "18 月+", "持续演进\n跨客户知识迁移\n世界模型探索", "L3-L4", ORANGE),
]
for i, (phase, name, time, desc, level, clr) in enumerate(roadmap):
    x = 0.8 + i * 3.1
    card(s, x, 1.6, 2.8, 2.5)
    bar(s, x, 1.6, 2.8, 0.08, clr)
    txt(s, x + 0.2, 1.8, 1.2, 0.3, phase, 14, clr, True)
    txt(s, x + 1.5, 1.8, 1.1, 0.3, time, 13, GRAY)
    txt(s, x + 0.2, 2.15, 2.4, 0.35, name, 18, WHITE, True)
    txt(s, x + 0.2, 2.55, 2.4, 1.0, desc, 14, LIGHT)
    txt(s, x + 0.2, 3.6, 2.4, 0.3, f"知识层级: {level}", 12, clr)

# 5-Day POV
txt(s, 0.8, 4.4, 12, 0.4, "5 天 POV -- 用客户数据证明价值", 20, WHITE, True)
pov_days = [
    ("Day 1", "环境部署 + 接入\n1-2 个核心数据源"),
    ("Day 2", "安装行业 AppPack\n+ Ontology 建模"),
    ("Day 3", "端到端演示\n自然语言->报告"),
    ("Day 4", "客户自行试用\n收集反馈"),
    ("Day 5", "价值对齐\nPOV 报告 + 商务"),
]
for i, (day, desc) in enumerate(pov_days):
    x = 0.8 + i * 2.46
    card(s, x, 4.9, 2.2, 1.3)
    txt(s, x + 0.15, 4.95, 1.9, 0.3, day, 16, ACCENT, True)
    txt(s, x + 0.15, 5.3, 1.9, 0.8, desc, 14, LIGHT)

# Core promises
card(s, 0.8, 6.4, 11.7, 0.9, DARK_GREEN)
txt(s, 1.2, 6.5, 11, 0.7,
    "核心承诺：不搬数据 (JDBC 直连) | 不改系统 (与现有 IT 共存) | 不绑架客户 (任何阶段可无损退出)\n"
    "「你之前的投资解决了'有没有数据'。DTS 解决'数据能不能做决策'。不是替换，是让已有投资增值。」", 15, LIGHT)


# ================================================================
# Slide 15: Competitive Landscape
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.8, 0.4, 12, 0.7, "竞争格局", 36, WHITE, True)

competitors = [
    ("Palantir\nFoundry/AIP", "产品成熟\n标杆客户(CIA/Airbus)\n重咨询 $1M+/年",
     "开放生态 vs 封闭\n中国市场\n价格 1/3-1/5", GRAY),
    ("Manus AI\n通用 Agent", "$100M ARR/8月\nMeta $2B 收购意向\n通用 Agent 执行器",
     "无行业深度\n无数据安全\n无知识沉淀", YELLOW),
    ("Coze/Dify\nAI 搭建平台", "灵活、上手快\n开发者友好",
     "企业级安全\n+ Ontology\n+ 行业知识", GRAY),
    ("阿里/华为\n大厂平台", "成熟数据治理\n大客户关系",
     "决策能力 vs\n只有数据工程", GRAY),
    ("DTS\n决策孪生系统", "连数据+沉淀知识\n+AI 决策+安全合规\n+行业能力包",
     "唯一同时具备\n五个要素的平台", GREEN),
]
for i, (name, desc, pos, clr) in enumerate(competitors):
    x = 0.8 + i * 2.46
    c_clr = DARK_GREEN if i == 4 else CARD
    card(s, x, 1.2, 2.2, 3.5, c_clr)
    txt(s, x + 0.1, 1.3, 2.0, 0.7, name, 14, WHITE, True)
    txt(s, x + 0.1, 2.1, 2.0, 1.3, desc, 13, LIGHT)
    txt(s, x + 0.1, 3.5, 2.0, 0.9, pos, 12, clr)

# Manus warning
card(s, 0.8, 5.0, 11.7, 1.1, DARK_RED)
txt(s, 1.1, 5.1, 11, 0.4,
    "Manus 的启示与警示", 18, ORANGE, True)
txt(s, 1.1, 5.5, 11, 0.5,
    "正面: Agent 市场需求真实  |  反面: 通用 Agent 无行业壁垒 -> 被模型厂商吞掉\n"
    "DTS 做「模型厂商不愿做的脏活累活」: 连数据、懂行业、管安全、沉淀知识", 15, LIGHT)

# Moat depth
card(s, 0.8, 6.3, 11.7, 1.0, DARK_BLUE)
txt(s, 1.2, 6.35, 11, 0.4,
    "护城河深度: DAP 协议生态 (3-5年) > KnowHow 资产锁定 (不可复制) > 行业 Ontology (2-3年) > 技术架构 (1-2年)", 15, LIGHT)


# ================================================================
# Slide 16: Five Iron Laws
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.8, 0.4, 12, 0.7, "五条不可动摇的原则", 36, WHITE, True)
txt(s, 0.8, 1.0, 12, 0.4, "客户最关心的问题：「万一 AI 出错怎么办？」-- 这五条就是答案", 16, GRAY)

laws = [
    ("1", "人工随时可接管", "删掉 AI，系统仍 100% 可用。任何步骤可人工接管", ACCENT),
    ("2", "核心安全不可绕过", "所有请求必经 dts-gateway，无旁路直连，AppPack 无特权", ACCENT2),
    ("3", "数据安全是基因", "所有数据出口必经 dts-data-security，AI 只看授权数据", GREEN),
    ("4", "全操作可追溯", "人和 AI 的操作 -> Kafka -> dts-audit-log，append-only 不可篡改", ORANGE),
    ("5", "能力先于界面", "API-first，先有完整 API + 测试，再做前端", PURPLE),
]
for i, (num, title, desc, clr) in enumerate(laws):
    y = 1.55 + i * 1.1
    card(s, 0.8, y, 11.7, 0.95)
    circle(s, 1.1, y + 0.15, 0.65, clr)
    txt(s, 1.1, y + 0.2, 0.65, 0.55, num, 24, WHITE, True, PP_ALIGN.CENTER)
    txt(s, 2.1, y + 0.1, 3.0, 0.4, title, 22, WHITE, True)
    txt(s, 2.1, y + 0.5, 9.5, 0.4, desc, 15, LIGHT)


# ================================================================
# Slide 17: Moats & Risks
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.8, 0.4, 12, 0.7, "竞争壁垒与风险", 36, WHITE, True)

moats = [
    ("DAP 协议生态", "协议+生态不可快速复制 (3-5 年)", ACCENT),
    ("KnowHow 资产锁定", "客户积累的行业知识不可复制", ACCENT2),
    ("行业 Ontology 深度", "需真实项目积累 (2-3 年)", GREEN),
    ("五条铁律架构保障", "架构级别，不是功能补丁 (1-2 年)", ORANGE),
    ("AppPack 网络效应", "开发者+客户同时积累 (2-3 年)", PURPLE),
]
for i, (n, d, c) in enumerate(moats):
    y = 1.2 + i * 0.68
    card(s, 0.8, y, 11.7, 0.58)
    txt(s, 1.1, y + 0.05, 3.2, 0.35, n, 16, c, True)
    txt(s, 4.3, y + 0.05, 8.0, 0.45, d, 14, LIGHT)

txt(s, 0.8, 4.7, 12, 0.45, "风险与应对", 22, WHITE, True)
risks = [
    ("LLM 幻觉/出错", "DAP 安全多层校验\n+ HITL + 铁律 #1"),
    ("大厂入局", "行业 KnowHow 壁垒\n大厂不做脏活累活"),
    ("客户不付费", "5 天 POV 真实数据\n用结果说话再签约"),
    ("通用 Agent 降维", "Manus 证明: 通用无壁垒\n行业深度才是差异化"),
]
for i, (risk, resp) in enumerate(risks):
    x = 0.8 + i * 3.1
    card(s, x, 5.2, 2.8, 1.4)
    txt(s, x + 0.2, 5.25, 2.4, 0.3, risk, 15, RED, True)
    txt(s, x + 0.2, 5.6, 2.4, 0.8, resp, 14, LIGHT)

# People risk
card(s, 0.8, 6.8, 11.7, 0.6)
txt(s, 1.1, 6.85, 3, 0.35, "人才风险应对", 15, YELLOW, True)
txt(s, 4.0, 6.85, 8, 0.45,
    "分离关注点: AI 工程师不需懂行业, 行业顾问不需懂 AI | 知识沉淀在系统中, 不在个人脑中", 14, LIGHT)


# ================================================================
# Slide 18: Financial Projections
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.8, 0.4, 12, 0.7, "财务预测", 36, WHITE, True)

# Revenue table
headers = ["", "Year 1", "Year 2", "Year 3"]
rows_fin = [
    ("平台基础费", "400 万", "1,500 万", "3,800 万"),
    ("AppPack 市场", "100 万", "700 万", "2,800 万"),
    ("KnowHow 交易", "0", "100 万", "900 万"),
    ("专业服务", "250 万", "800 万", "2,000 万"),
    ("硬件集成", "50 万", "400 万", "500 万"),
    ("总收入", "800 万", "3,500 万", "1 亿"),
]
cw = [2.5, 2.8, 2.8, 2.8]
for ri, rd in enumerate([-1] + list(range(len(rows_fin)))):
    y = 1.2 + ri * 0.48
    if rd == -1:
        x_pos = 0.8
        for j, h in enumerate(headers):
            txt(s, x_pos, y, cw[j], 0.40, h, 15, ACCENT, True); x_pos += cw[j]
    else:
        d = rows_fin[rd]; total = d[0] == "总收入"; x_pos = 0.8
        if total:
            bar(s, 0.8, y, sum(cw), 0.40, DARK_BLUE)
        for j, v in enumerate(d):
            c = GREEN if total else (GRAY if j == 0 else LIGHT)
            txt(s, x_pos, y, cw[j], 0.40, v, 16 if not total else 18, c, total); x_pos += cw[j]

# Key metrics
txt(s, 0.8, 4.8, 12, 0.45, "关键指标", 22, WHITE, True)
kmetrics = [
    ("客户数", "6", "18", "48"),
    ("续费率", "-", "85%", "90%"),
    ("EBITDA", "-200万", "1,700万", "7,100万"),
    ("利润率", "-25%", "49%", "71%"),
]
for i, (name, y1, y2, y3) in enumerate(kmetrics):
    x = 0.8 + i * 3.1
    card(s, x, 5.3, 2.8, 1.8)
    txt(s, x + 0.2, 5.35, 2.4, 0.3, name, 16, GRAY)
    txt(s, x + 0.2, 5.7, 0.8, 0.4, y1, 14, LIGHT)
    txt(s, x + 1.0, 5.7, 0.8, 0.4, y2, 14, LIGHT)
    txt(s, x + 1.8, 5.7, 0.8, 0.4, y3, 18, GREEN, True)
    txt(s, x + 0.2, 6.1, 0.8, 0.25, "Y1", 11, GRAY)
    txt(s, x + 1.0, 6.1, 0.8, 0.25, "Y2", 11, GRAY)
    txt(s, x + 1.8, 6.1, 0.8, 0.25, "Y3", 11, GRAY)

txt(s, 0.8, 7.15, 12, 0.3,
    "Y1 亏损可控 (200 万) | 盈亏平衡: 第 15 个月 | Y3 利润率 71% -- 知识资产平台的典型曲线", 15, GRAY)


# ================================================================
# Slide 19: Funding Ask
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.8, 0.4, 12, 0.7, "融资计划", 36, WHITE, True)

# Key terms
card(s, 0.8, 1.3, 5.5, 2.5, DARK_BLUE)
txt(s, 1.2, 1.5, 5.0, 0.5, "种子轮", 28, ACCENT, True)
terms = [
    ("融资金额", "2,000 万 RMB"),
    ("出让股份", "10-15%"),
    ("投前估值", "1.3-2.0 亿 RMB"),
    ("估值逻辑", "Y3 收入 1 亿 x 8-12x PS -> 种子轮折扣 60-80%"),
]
for i, (k, v) in enumerate(terms):
    y = 2.1 + i * 0.4
    txt(s, 1.2, y, 1.5, 0.3, k, 14, GRAY)
    txt(s, 2.8, y, 3.5, 0.3, v, 15, WHITE, True)

# Fund allocation
card(s, 7.0, 1.3, 5.5, 2.5)
txt(s, 7.3, 1.5, 5.0, 0.4, "资金用途", 20, WHITE, True)
allocs = [
    ("核心研发", "60%", "1,200 万", ACCENT),
    ("市场拓展", "25%", "500 万", GREEN),
    ("运营储备", "15%", "300 万", ORANGE),
]
for i, (name, pct, amt, clr) in enumerate(allocs):
    y = 2.0 + i * 0.55
    bar(s, 7.3, y, 4.8 * (int(pct[:-1]) / 60), 0.35, clr)
    txt(s, 7.5, y + 0.02, 2, 0.3, f"{name} ({pct})", 14, WHITE, True)
    txt(s, 10.5, y + 0.02, 1.5, 0.3, amt, 14, WHITE)

# Milestones
txt(s, 0.8, 4.1, 12, 0.4, "里程碑", 20, WHITE, True)
milestones = [
    ("M1: MVP", "+6 月", "25 微服务上线\n端到端跑通", ACCENT),
    ("M2: 首个标杆", "+9 月", "1 个付费客户\n年合同 100 万+", ACCENT2),
    ("M3: Pack 验证", "+12 月", "2 个 AppPack\n复用率 > 50%", GREEN),
    ("M4: 准备 A 轮", "+15 月", "ARR 500 万+\n5+ 付费客户", ORANGE),
]
for i, (name, time, desc, clr) in enumerate(milestones):
    x = 0.8 + i * 3.1
    card(s, x, 4.5, 2.8, 1.5)
    bar(s, x, 4.5, 2.8, 0.06, clr)
    txt(s, x + 0.2, 4.65, 1.5, 0.3, name, 15, clr, True)
    txt(s, x + 1.8, 4.65, 0.8, 0.3, time, 13, GRAY)
    txt(s, x + 0.2, 5.0, 2.4, 0.8, desc, 14, LIGHT)

# Follow-on rounds
txt(s, 0.8, 6.2, 12, 0.4, "后续融资 & 退出路径", 20, WHITE, True)
card(s, 0.8, 6.6, 5.5, 0.8)
txt(s, 1.1, 6.65, 5.0, 0.3,
    "A 轮 (+15-18月): 5,000-8,000 万 | B 轮 (+30月): 2-3 亿", 14, LIGHT)

card(s, 7.0, 6.6, 5.5, 0.8)
txt(s, 7.3, 6.65, 5.0, 0.3,
    "IPO (科创板) | 战略并购 (华为/阿里) | 生态拆分上市", 14, LIGHT)


# ================================================================
# Slide 20: Summary
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
bar(s, 1, 1.8, 1.5, 0.06, ACCENT)
txt(s, 1, 2.1, 11, 1.0, "Decision Twins System", 44, WHITE, True)
txt(s, 1, 3.0, 11, 0.8, "把专家脑中的隐性经验变成系统里能自动运转的能力包", 26, ACCENT)

bullets(s, 1, 4.0, 11, 2.5, [
    "核心公式: Ontology (认知基座) + Intent (意图理解) + Agent (自主执行)",
    "三大支柱: dts-studio (设计) + dts-stack (运行) + app-stack (应用)",
    "25 微服务 + 57 AI 技能 + DAP 五层协议构建生态护城河",
    "知识战略: Polanyi 缄默知识 + LeCun 1% 领域数据涌现",
    "商业模式: 平台基础费 + AppPack 市场 + KnowHow 交易 -- 越用越值钱",
    "五条铁律: 人工接管 | 安全不绕过 | 数据安全基因 | 全操作追溯 | 能力先于界面",
    "Y3 目标: 1 亿收入, 71% 利润率, 48 个客户",
], 19, LIGHT)

txt(s, 1, 6.5, 11, 0.5,
    "种子轮融资 2,000 万 -- 邀请您共同定义企业 AI 决策的未来", 22, GREEN, True)
txt(s, 1, 7.0, 11, 0.4,
    "Ontology + Intent + Agent = Decision Twins", 16, GRAY, False, PP_ALIGN.LEFT)


# -- Save --
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "DTS-商业计划书-V1.0.0.pptx")
prs.save(out)
print(f"Saved: {out}")
