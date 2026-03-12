#!/usr/bin/env python3
"""Generate DTS Deployment & Hardware Pricing PPT — V1.0.0 edition."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# -- Colors --
BG = RGBColor(0x1A, 0x1A, 0x2E)
CARD = RGBColor(0x22, 0x22, 0x3A)
ACCENT = RGBColor(0x00, 0x96, 0xD6)
ACCENT2 = RGBColor(0x00, 0xC9, 0xA7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xAA, 0xAA, 0xBB)
LIGHT = RGBColor(0xE0, 0xE0, 0xE8)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
RED = RGBColor(0xE8, 0x4D, 0x4D)
GREEN = RGBColor(0x28, 0xC7, 0x6F)
PURPLE = RGBColor(0xBB, 0x86, 0xFC)
DARK_BLUE = RGBColor(0x0A, 0x3D, 0x5C)
DARK_RED = RGBColor(0x3A, 0x1A, 0x1A)
DARK_GREEN = RGBColor(0x0A, 0x3D, 0x2A)
DARK_PURPLE = RGBColor(0x2A, 0x1A, 0x3A)
YELLOW = RGBColor(0xFF, 0xD7, 0x00)


def bg(slide):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = BG

def txt(slide, l, t, w, h, text, sz=18, clr=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = clr; p.font.bold = bold
    p.font.name = "Microsoft YaHei"; p.alignment = align
    return tb

def card(slide, l, t, w, h, clr=CARD):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = clr; s.line.fill.background(); s.shadow.inherit = False
    return s

def bullets(slide, l, t, w, h, items, sz=16, clr=LIGHT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(sz); p.font.color.rgb = clr
        p.font.name = "Microsoft YaHei"; p.space_after = Pt(8)
    return tb

def bar(slide, l, t, w, h, clr):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = clr; s.line.fill.background()
    return s

def circle(slide, l, t, d, clr):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(d), Inches(d))
    s.fill.solid(); s.fill.fore_color.rgb = clr; s.line.fill.background()
    return s


# ================================================================
# Slide 1: Cover
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
bar(s, 1, 2.3, 1.5, 0.06, ACCENT)
txt(s, 1, 2.6, 11, 1.2, "DTS 部署模式与硬件定价方案", 48, WHITE, True)
txt(s, 1, 3.7, 11, 0.8, "Decision Twins System — Deployment & Pricing V1.0.0", 28, ACCENT)
txt(s, 1, 4.6, 11, 0.6, "软硬一体交付，覆盖中小企业到政府信创全场景", 20, GRAY)
txt(s, 1, 5.4, 8, 0.5, "2026-03  |  CONFIDENTIAL", 16, GRAY)


# ================================================================
# Slide 2: Market Insight
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
bar(s, 0.6, 0.4, 0.08, 0.5, ACCENT)
txt(s, 0.8, 0.4, 12, 0.7, "市场洞察 — 为什么需要软硬一体", 36, WHITE, True)

insights = [
    ("软件无形，硬件有形", "客户更愿意为看得见摸得着的硬件付费\n软件价值认知薄弱", ACCENT, DARK_BLUE),
    ("预算科目限制", "很多企事业单位有\"设备采购\"预算\n但没有\"软件采购\"预算", ORANGE, DARK_RED),
    ("国产化/信创要求", "政府、央企、军工必须采购\n国产软硬件，通过信创认证", GREEN, DARK_GREEN),
]
for i, (title, desc, clr, bg_clr) in enumerate(insights):
    x = 0.6 + i * 4.1
    card(s, x, 1.5, 3.8, 3.5, bg_clr)
    circle(s, x + 0.3, 1.8, 0.7, clr)
    txt(s, x + 0.3, 1.9, 0.7, 0.5, str(i + 1), 24, WHITE, True, PP_ALIGN.CENTER)
    txt(s, x + 1.2, 1.85, 2.5, 0.5, title, 22, clr, True)
    txt(s, x + 0.3, 2.8, 3.3, 2.0, desc, 16, LIGHT)

card(s, 0.6, 5.4, 12.1, 1.2, DARK_BLUE)
txt(s, 1.0, 5.5, 4, 0.5, "DTS 的回答：两种交付模式", 22, ACCENT, True)
txt(s, 1.0, 6.0, 5, 0.5, "纯软件 — 客户自备硬件/云资源，DTS 提供软件授权", 16, LIGHT)
txt(s, 7.0, 6.0, 5.5, 0.5, "软硬一体 — DTS 预装在认证硬件上，整机交付，开箱即用", 16, ACCENT2)


# ================================================================
# Slide 3: Design Principles
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
bar(s, 0.6, 0.4, 0.08, 0.5, ACCENT)
txt(s, 0.8, 0.4, 12, 0.7, "设计原则", 36, WHITE, True)

principles = [
    ("01", "开箱即用", "通电联网即可使用\nDay 1 部署 < 1 小时", ACCENT),
    ("02", "平滑升级", "一体机 → 集群\n迁移路径清晰，数据零丢失", GREEN),
    ("03", "灵活定价", "买断/订阅/整机打包\n适配不同投标场景", ORANGE),
    ("04", "信创兼容", "国产 CPU + GPU + OS\n信创目录认证", PURPLE),
    ("05", "硬件开放", "不绑定特定厂商\nBOM 清单认证制", ACCENT2),
]
for i, (num, title, desc, clr) in enumerate(principles):
    x = 0.6 + i * 2.5
    card(s, x, 1.5, 2.2, 4.5, CARD)
    bar(s, x, 1.5, 2.2, 0.06, clr)
    txt(s, x + 0.2, 1.8, 1.8, 0.5, num, 36, clr, True)
    txt(s, x + 0.2, 2.4, 1.8, 0.5, title, 22, WHITE, True)
    txt(s, x + 0.2, 3.1, 1.8, 2.5, desc, 15, LIGHT)


# ================================================================
# Slide 4: Product Overview — Two Modes
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
bar(s, 0.6, 0.4, 0.08, 0.5, ACCENT)
txt(s, 0.8, 0.4, 12, 0.7, "两大交付模式总览", 36, WHITE, True)

# All-in-one box
card(s, 0.6, 1.5, 5.8, 5.0, DARK_BLUE)
txt(s, 1.0, 1.7, 5, 0.5, "DTS-Box 一体机", 28, ACCENT, True)
bar(s, 1.0, 2.3, 4.5, 0.04, ACCENT)

box_items = [
    "单台服务器承载 DTS 全栈",
    "塔式（中小企业）/ 机架式（政府信创）",
    "适用 5-50 人团队",
    "无需专职 IT 运维",
    "支持离线/涉密网络",
    "预算 ≤ 50 万",
    "三档 SKU: Lite / Pro / Gov",
]
bullets(s, 1.0, 2.6, 5, 3.5, box_items, 16, LIGHT)

# Cluster
card(s, 6.9, 1.5, 5.8, 5.0, DARK_GREEN)
txt(s, 7.3, 1.7, 5, 0.5, "DTS-Cluster 三节点集群", 28, GREEN, True)
bar(s, 7.3, 2.3, 4.5, 0.04, GREEN)

cluster_items = [
    "三台服务器组成 HA 集群",
    "2 台业务互备 + 1 台独立 GPU",
    "适用 100+ 并发用户",
    "生产环境不允许停机",
    "TB 级数据规模",
    "预算 100-400 万",
    "两档 SKU: Standard / Gov",
]
bullets(s, 7.3, 2.6, 5, 3.5, cluster_items, 16, LIGHT)


# ================================================================
# Slide 5: DTS-Box Lite
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
bar(s, 0.6, 0.4, 0.08, 0.5, ACCENT)
txt(s, 0.8, 0.4, 12, 0.7, "DTS-Box Lite — 入门款", 36, WHITE, True)
txt(s, 9.5, 0.4, 3.5, 0.7, "¥10万", 42, ACCENT, True, PP_ALIGN.RIGHT)

# Hardware specs
card(s, 0.6, 1.5, 5.8, 4.0, CARD)
txt(s, 1.0, 1.6, 5, 0.5, "硬件配置", 22, ACCENT, True)
bar(s, 1.0, 2.1, 4.5, 0.04, ACCENT)

specs = [
    ("CPU", "8C16T+ (x86/ARM)"),
    ("内存", "64 GB"),
    ("存储", "1TB NVMe SSD"),
    ("GPU", "无（纯 CPU 推理或外接 API）"),
    ("K8s", "K3s 单节点"),
]
for i, (k, v) in enumerate(specs):
    y = 2.4 + i * 0.45
    txt(s, 1.0, y, 1.5, 0.4, k, 16, GRAY)
    txt(s, 2.5, y, 3.5, 0.4, v, 16, WHITE)

# Capabilities
card(s, 6.9, 1.5, 5.8, 4.0, CARD)
txt(s, 7.3, 1.6, 5, 0.5, "能力与场景", 22, GREEN, True)
bar(s, 7.3, 2.1, 4.5, 0.04, GREEN)

cap_items = [
    "模型能力: 7B 量化模型 / 云端 LLM",
    "并发用户: 5-15 人",
    "DTS 全栈 25 个微服务",
    "全部中间件内置",
    "",
    "典型场景:",
    "  预算有限的中小企业试点",
    "  边缘节点接入中心集群",
]
bullets(s, 7.3, 2.3, 5, 3.0, cap_items, 15, LIGHT)

# Bottom note
card(s, 0.6, 5.8, 12.1, 1.0, DARK_BLUE)
txt(s, 1.0, 5.9, 11, 0.7, "形态: 塔式工作站（办公室）或 1U 机架式（机房）  |  含首年软件授权 + 实施部署 + 维保", 18, LIGHT)


# ================================================================
# Slide 6: DTS-Box Pro
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
bar(s, 0.6, 0.4, 0.08, 0.5, ACCENT)
txt(s, 0.8, 0.4, 12, 0.7, "DTS-Box Pro — 标准款", 36, WHITE, True)
txt(s, 9.5, 0.4, 3.5, 0.7, "¥30万", 42, ORANGE, True, PP_ALIGN.RIGHT)

card(s, 0.6, 1.5, 5.8, 4.0, CARD)
txt(s, 1.0, 1.6, 5, 0.5, "硬件配置", 22, ORANGE, True)
bar(s, 1.0, 2.1, 4.5, 0.04, ORANGE)

specs = [
    ("CPU", "16C32T+ (x86/ARM)"),
    ("内存", "128 GB"),
    ("存储", "2TB NVMe SSD + 4TB HDD"),
    ("GPU", "NVIDIA RTX 4090 24GB 或同级"),
    ("K8s", "K3s 单节点，中间件全内置"),
]
for i, (k, v) in enumerate(specs):
    y = 2.4 + i * 0.45
    txt(s, 1.0, y, 1.5, 0.4, k, 16, GRAY)
    txt(s, 2.5, y, 3.5, 0.4, v, 16, WHITE)

card(s, 6.9, 1.5, 5.8, 4.0, CARD)
txt(s, 7.3, 1.6, 5, 0.5, "能力与场景", 22, ORANGE, True)
bar(s, 7.3, 2.1, 4.5, 0.04, ORANGE)

cap_items = [
    "模型能力: 14B-32B，推理速度良好",
    "并发用户: 15-50 人",
    "完整离线运行 DTS 全栈",
    "本地 LLM 推理（Ollama/vLLM）",
    "",
    "典型场景:",
    "  中小企业主力款",
    "  离线环境完整运行",
]
bullets(s, 7.3, 2.3, 5, 3.0, cap_items, 15, LIGHT)

card(s, 0.6, 5.8, 12.1, 1.0, DARK_GREEN)
txt(s, 1.0, 5.9, 11, 0.7, "中小企业最佳性价比之选  |  含首年软件授权 + 实施部署 + 维保 + 基础 AppPack", 18, ACCENT2)


# ================================================================
# Slide 7: DTS-Box Gov
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
bar(s, 0.6, 0.4, 0.08, 0.5, ACCENT)
txt(s, 0.8, 0.4, 12, 0.7, "DTS-Box Gov — 信创款", 36, WHITE, True)
txt(s, 9.5, 0.4, 3.5, 0.7, "¥45万", 42, PURPLE, True, PP_ALIGN.RIGHT)

card(s, 0.6, 1.5, 5.8, 4.5, CARD)
txt(s, 1.0, 1.6, 5, 0.5, "硬件配置（全国产化）", 22, PURPLE, True)
bar(s, 1.0, 2.1, 4.5, 0.04, PURPLE)

specs = [
    ("CPU", "鲲鹏 920 / 飞腾 S2500"),
    ("内存", "128-256 GB"),
    ("存储", "2TB 国产 NVMe SSD"),
    ("GPU", "昇腾 310B / 910B"),
    ("OS", "麒麟 V10 / 统信 UOS"),
    ("K8s", "K3s (ARM64) 单节点"),
    ("认证", "信创目录认证"),
]
for i, (k, v) in enumerate(specs):
    y = 2.4 + i * 0.42
    txt(s, 1.0, y, 1.5, 0.4, k, 16, GRAY)
    txt(s, 2.5, y, 3.5, 0.4, v, 16, WHITE)

card(s, 6.9, 1.5, 5.8, 4.5, CARD)
txt(s, 7.3, 1.6, 5, 0.5, "能力与场景", 22, PURPLE, True)
bar(s, 7.3, 2.1, 4.5, 0.04, PURPLE)

cap_items = [
    "模型能力: 14B-70B（视昇腾型号）",
    "并发用户: 15-50 人",
    "全栈国产化替代",
    "物理隔离/涉密网络支持",
    "",
    "典型场景:",
    "  政府/央企/军工",
    "  满足国产化替代硬性要求",
    "  信创目录认证项目",
]
bullets(s, 7.3, 2.3, 5, 3.5, cap_items, 15, LIGHT)

card(s, 0.6, 6.3, 12.1, 0.7, DARK_PURPLE)
txt(s, 1.0, 6.35, 11, 0.5, "全栈信创  |  国产 CPU + GPU + OS  |  满足等保/分保要求  |  含首年授权 + 实施 + 维保", 18, LIGHT)


# ================================================================
# Slide 8: Box SKU Comparison
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
bar(s, 0.6, 0.4, 0.08, 0.5, ACCENT)
txt(s, 0.8, 0.4, 12, 0.7, "一体机三档 SKU 对比", 36, WHITE, True)

# Table header
cols = [("", 2.0), ("Lite", 3.2), ("Pro", 3.2), ("Gov", 3.2)]
x_positions = [0.6, 2.6, 5.8, 9.0]
for i, ((label, w), x) in enumerate(zip(cols, x_positions)):
    if i > 0:
        card(s, x, 1.3, w, 0.6, DARK_BLUE if i == 1 else (DARK_RED if i == 2 else DARK_PURPLE))
        clr = ACCENT if i == 1 else (ORANGE if i == 2 else PURPLE)
        txt(s, x + 0.1, 1.35, w - 0.2, 0.5, label, 22, clr, True, PP_ALIGN.CENTER)

rows = [
    ("报价", "¥10万", "¥30万", "¥45万"),
    ("CPU", "8C16T+", "16C32T+", "鲲鹏/飞腾"),
    ("内存", "64 GB", "128 GB", "128-256 GB"),
    ("GPU", "无", "RTX 4090 24GB", "昇腾 310B/910B"),
    ("模型", "7B / 云端 API", "14B-32B", "14B-70B"),
    ("并发", "5-15 人", "15-50 人", "15-50 人"),
    ("信创", "—", "—", "全栈认证"),
    ("场景", "试点/边缘", "中小企业主力", "政府/军工"),
]
for r, (label, v1, v2, v3) in enumerate(rows):
    y = 2.1 + r * 0.55
    if r % 2 == 0:
        bar(s, 0.6, y, 11.6, 0.5, RGBColor(0x1E, 0x1E, 0x32))
    txt(s, 0.7, y + 0.05, 1.8, 0.4, label, 16, GRAY, True)
    txt(s, 2.7, y + 0.05, 3.0, 0.4, v1, 16, WHITE, False, PP_ALIGN.CENTER)
    txt(s, 5.9, y + 0.05, 3.0, 0.4, v2, 16, WHITE, False, PP_ALIGN.CENTER)
    txt(s, 9.1, y + 0.05, 3.0, 0.4, v3, 16, WHITE, False, PP_ALIGN.CENTER)

# Common features
card(s, 0.6, 6.5, 12.1, 0.7, CARD)
txt(s, 1.0, 6.55, 11.5, 0.5, "共性: 全部预装 DTS 全栈 25 微服务 + 全部中间件 + 系统初始化向导 + 含首年授权与维保", 16, ACCENT2)


# ================================================================
# Slide 9: Software Pre-install
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
bar(s, 0.6, 0.4, 0.08, 0.5, ACCENT)
txt(s, 0.8, 0.4, 12, 0.7, "出厂预装软件清单", 36, WHITE, True)

categories = [
    ("DTS 平台", ACCENT, [
        "Java 业务平台 (10 svc)",
        "Python AI 核心 (10 svc)",
        "Go 基础设施 (2 svc)",
        "Frontend (3 webapps)",
    ]),
    ("中间件", GREEN, [
        "PostgreSQL (数据存储)",
        "Kafka KRaft (消息总线)",
        "ClickHouse (分析引擎)",
        "Neo4j CE (图数据库)",
        "pgvector + MinIO",
        "Keycloak (认证)",
    ]),
    ("AI 运行时", ORANGE, [
        "预置本地 LLM (按档位)",
        "Ollama / vLLM 推理",
        "预置 Embedding 模型",
        "模型管理工具",
    ]),
    ("管理工具", PURPLE, [
        "dts-cli (命令行)",
        "Grafana (可观测)",
        "系统初始化向导",
        "数据备份/恢复",
    ]),
]
for i, (title, clr, items) in enumerate(categories):
    x = 0.6 + i * 3.15
    card(s, x, 1.4, 2.9, 5.5, CARD)
    bar(s, x, 1.4, 2.9, 0.06, clr)
    txt(s, x + 0.2, 1.6, 2.5, 0.5, title, 22, clr, True)
    bullets(s, x + 0.2, 2.3, 2.5, 4.0, items, 15, LIGHT)


# ================================================================
# Slide 10: Cluster Architecture
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
bar(s, 0.6, 0.4, 0.08, 0.5, ACCENT)
txt(s, 0.8, 0.4, 12, 0.7, "三节点集群架构", 36, WHITE, True)

# Outer frame
card(s, 0.6, 1.3, 12.1, 5.8, RGBColor(0x18, 0x18, 0x28))
txt(s, 1.0, 1.4, 5, 0.5, "K8s HA Cluster (etcd 3 副本)", 18, ACCENT, True)

# Node 1
card(s, 0.9, 2.1, 3.8, 4.5, DARK_BLUE)
txt(s, 1.2, 2.2, 3.4, 0.5, "Node 1 — 业务+管理", 20, ACCENT, True)
bar(s, 1.2, 2.7, 3.0, 0.04, ACCENT)
items_n1 = ["Java 10 services", "Go 2 services", "Frontend 3 apps", "中间件 (HA 主)", "etcd member"]
bullets(s, 1.2, 2.9, 3.2, 3.5, items_n1, 14, LIGHT)

# Node 2
card(s, 5.0, 2.1, 3.8, 4.5, DARK_BLUE)
txt(s, 5.3, 2.2, 3.4, 0.5, "Node 2 — 业务+管理", 20, ACCENT, True)
bar(s, 5.3, 2.7, 3.0, 0.04, ACCENT)
items_n2 = ["Java 10 services", "Go 2 services", "Frontend 3 apps", "中间件 (HA 备)", "etcd member"]
bullets(s, 5.3, 2.9, 3.2, 3.5, items_n2, 14, LIGHT)

# Mutual backup arrow text
txt(s, 3.4, 4.8, 3.2, 0.5, "互为备份", 16, ACCENT2, True, PP_ALIGN.CENTER)

# Node 3 — GPU
card(s, 9.1, 2.1, 3.4, 4.5, DARK_GREEN)
txt(s, 9.4, 2.2, 3.0, 0.5, "Node 3 — GPU", 20, GREEN, True)
bar(s, 9.4, 2.7, 2.5, 0.04, GREEN)
items_n3 = ["Python AI 10 services", "LLM 推理引擎", "vLLM / Ollama", "GPU: L40S/A800", "或 昇腾 910B", "etcd member"]
bullets(s, 9.4, 2.9, 2.8, 3.5, items_n3, 14, LIGHT)

txt(s, 8.2, 4.8, 2.0, 0.5, "独立算力", 16, GREEN, True, PP_ALIGN.CENTER)


# ================================================================
# Slide 11: Cluster HA Features
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
bar(s, 0.6, 0.4, 0.08, 0.5, ACCENT)
txt(s, 0.8, 0.4, 12, 0.7, "集群高可用特性", 36, WHITE, True)

ha_features = [
    ("K8s HA", "etcd 三副本\n任一节点故障自动选主\nPod 自动调度迁移", ACCENT, DARK_BLUE),
    ("中间件 HA", "PG 主备自动切换\nKafka 多 Broker\nClickHouse 分片副本", GREEN, DARK_GREEN),
    ("GPU 独立", "算力节点独立升级\n不影响业务连续性\n按需扩展 GPU 型号", ORANGE, DARK_RED),
    ("弹性扩展", "保留向 5-10 节点\n扩展的迁移路径\n水平扩容无中断", PURPLE, DARK_PURPLE),
]
for i, (title, desc, clr, bg_clr) in enumerate(ha_features):
    x = 0.6 + i * 3.15
    card(s, x, 1.4, 2.9, 4.8, bg_clr)
    circle(s, x + 1.0, 1.7, 0.9, clr)
    txt(s, x + 1.0, 1.8, 0.9, 0.7, "HA", 18, WHITE, True, PP_ALIGN.CENTER)
    txt(s, x + 0.2, 2.9, 2.5, 0.5, title, 22, clr, True)
    txt(s, x + 0.2, 3.5, 2.5, 2.5, desc, 16, LIGHT)


# ================================================================
# Slide 12: Cluster SKU
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
bar(s, 0.6, 0.4, 0.08, 0.5, ACCENT)
txt(s, 0.8, 0.4, 12, 0.7, "集群两档 SKU", 36, WHITE, True)

# Standard
card(s, 0.6, 1.3, 5.8, 5.5, DARK_BLUE)
txt(s, 1.0, 1.5, 5, 0.5, "DTS-Cluster Standard", 26, ACCENT, True)
txt(s, 1.0, 2.1, 5, 0.5, "¥100-200万", 32, WHITE, True)
bar(s, 1.0, 2.7, 4.5, 0.04, ACCENT)

std_specs = [
    "业务节点 ×2:",
    "  32C64T, 256GB RAM, 4TB NVMe",
    "",
    "GPU 节点 ×1:",
    "  32C64T, 256GB, L40S 48GB / A800 80GB",
    "",
    "并发 100-500 用户 | TB 级数据",
    "70B+ 模型高速推理",
]
bullets(s, 1.0, 2.9, 5, 3.5, std_specs, 15, LIGHT)

# Gov
card(s, 6.9, 1.3, 5.8, 5.5, DARK_PURPLE)
txt(s, 7.3, 1.5, 5, 0.5, "DTS-Cluster Gov", 26, PURPLE, True)
txt(s, 7.3, 2.1, 5, 0.5, "¥200-400万", 32, WHITE, True)
bar(s, 7.3, 2.7, 4.5, 0.04, PURPLE)

gov_specs = [
    "业务节点 ×2:",
    "  鲲鹏920 32C, 256GB, 4TB 国产 SSD",
    "",
    "GPU 节点 ×1:",
    "  鲲鹏920 + 昇腾 910B, 256GB",
    "",
    "OS: 麒麟 V10 / 统信 UOS",
    "全栈信创目录认证",
]
bullets(s, 7.3, 2.9, 5, 3.5, gov_specs, 15, LIGHT)


# ================================================================
# Slide 13: Pricing Overview
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
bar(s, 0.6, 0.4, 0.08, 0.5, ACCENT)
txt(s, 0.8, 0.4, 12, 0.7, "全线产品价格总览", 36, WHITE, True)

skus = [
    ("DTS-Box\nLite", "¥10万", ACCENT, DARK_BLUE),
    ("DTS-Box\nPro", "¥30万", ORANGE, DARK_RED),
    ("DTS-Box\nGov", "¥45万", PURPLE, DARK_PURPLE),
    ("Cluster\nStandard", "¥100-200万", GREEN, DARK_GREEN),
    ("Cluster\nGov", "¥200-400万", RED, DARK_RED),
]
for i, (name, price, clr, bg_clr) in enumerate(skus):
    x = 0.6 + i * 2.5
    card(s, x, 1.4, 2.2, 3.5, bg_clr)
    bar(s, x, 1.4, 2.2, 0.06, clr)
    txt(s, x + 0.15, 1.7, 1.9, 1.0, name, 20, WHITE, True, PP_ALIGN.CENTER)
    txt(s, x + 0.15, 2.9, 1.9, 0.5, price, 24, clr, True, PP_ALIGN.CENTER)
    txt(s, x + 0.15, 3.5, 1.9, 1.0, "含首年服务", 14, GRAY, False, PP_ALIGN.CENTER)

# Scale bar visualization
txt(s, 0.8, 5.3, 12, 0.5, "价格带分布", 20, WHITE, True)
bar(s, 0.8, 5.8, 2.0, 0.4, ACCENT)  # Lite 10
txt(s, 0.8, 5.85, 2.0, 0.35, "10万", 14, WHITE, False, PP_ALIGN.CENTER)
bar(s, 3.0, 5.8, 3.0, 0.4, ORANGE)  # Pro 30
txt(s, 3.0, 5.85, 3.0, 0.35, "30万", 14, WHITE, False, PP_ALIGN.CENTER)
bar(s, 6.2, 5.8, 2.0, 0.4, PURPLE)  # Gov 45
txt(s, 6.2, 5.85, 2.0, 0.35, "45万", 14, WHITE, False, PP_ALIGN.CENTER)
bar(s, 8.4, 5.8, 2.0, 0.4, GREEN)  # Cluster Std
txt(s, 8.4, 5.85, 2.0, 0.35, "100-200万", 12, WHITE, False, PP_ALIGN.CENTER)
bar(s, 10.6, 5.8, 2.2, 0.4, RED)  # Cluster Gov
txt(s, 10.6, 5.85, 2.2, 0.35, "200-400万", 12, WHITE, False, PP_ALIGN.CENTER)

txt(s, 0.8, 6.4, 3, 0.4, "一体机 (≤50万)", 16, GRAY)
txt(s, 8.4, 6.4, 5, 0.4, "集群 (100-400万)", 16, GRAY)


# ================================================================
# Slide 14: Pricing Modes
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
bar(s, 0.6, 0.4, 0.08, 0.5, ACCENT)
txt(s, 0.8, 0.4, 12, 0.7, "三种报价方式 — 按投标场景灵活选择", 36, WHITE, True)

modes = [
    ("方式 A", "一次性买断", ACCENT, DARK_BLUE, [
        "硬件采购费",
        "软件永久授权",
        "首年实施 + 维保",
        "次年起维保费 (报价 15-20%)",
        "",
        "适合: 有设备采购预算",
        "买完即拥有，无后续绑定",
    ]),
    ("方式 B", "硬件+订阅", GREEN, DARK_GREEN, [
        "硬件按成本价（可补贴）",
        "软件年费订阅 (3-5 年锁定)",
        "首年实施费",
        "年费含维保 + 升级",
        "",
        "适合: 锁定长期收入",
        "客户初期投入低",
    ]),
    ("方式 C", "整机方案价", ORANGE, DARK_RED, [
        "不拆分硬件/软件",
        "按\"AI 决策解决方案\"报价",
        "利润内含，简化流程",
        "一张发票搞定",
        "",
        "适合: 预算科目为\"设备\"",
        "客户不关心软硬拆分",
    ]),
]
for i, (label, title, clr, bg_clr, items) in enumerate(modes):
    x = 0.6 + i * 4.1
    card(s, x, 1.3, 3.8, 5.5, bg_clr)
    bar(s, x, 1.3, 3.8, 0.06, clr)
    txt(s, x + 0.3, 1.6, 3.2, 0.5, label, 16, GRAY)
    txt(s, x + 0.3, 2.0, 3.2, 0.5, title, 26, clr, True)
    bullets(s, x + 0.3, 2.7, 3.2, 4.0, items, 15, LIGHT)


# ================================================================
# Slide 15: Cost Breakdown
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
bar(s, 0.6, 0.4, 0.08, 0.5, ACCENT)
txt(s, 0.8, 0.4, 12, 0.7, "价格构成拆解（内部核算）", 36, WHITE, True)

# Headers
headers = ["成本项", "Lite", "Pro", "Gov", "Cluster Std", "Cluster Gov"]
widths = [2.0, 1.8, 1.8, 1.8, 2.2, 2.2]
x_pos = [0.6]
for w in widths[:-1]:
    x_pos.append(x_pos[-1] + w + 0.1)

for i, (h, x, w) in enumerate(zip(headers, x_pos, widths)):
    card(s, x, 1.3, w, 0.6, DARK_BLUE if i == 0 else CARD)
    txt(s, x + 0.1, 1.35, w - 0.2, 0.5, h, 16, ACCENT if i == 0 else WHITE, True, PP_ALIGN.CENTER)

rows = [
    ("硬件成本", "2-3万", "8-12万", "15-20万", "30-50万", "60-100万"),
    ("软件授权", "3-4万", "10-12万", "12-15万", "40-80万", "80-160万"),
    ("实施部署", "1-2万", "3-5万", "5-8万", "15-30万", "30-60万"),
    ("首年维保", "含", "含", "含", "含", "含"),
    ("毛利率", "~40%", "~40%", "~35%", "~40%", "~40%"),
]
for r, row in enumerate(rows):
    y = 2.1 + r * 0.7
    if r % 2 == 0:
        bar(s, 0.6, y, 11.7, 0.6, RGBColor(0x1E, 0x1E, 0x32))
    for c, (val, x, w) in enumerate(zip(row, x_pos, widths)):
        clr = GRAY if c == 0 else (GREEN if r == 4 else WHITE)
        bold = c == 0 or r == 4
        txt(s, x + 0.1, y + 0.1, w - 0.2, 0.4, val, 15, clr, bold, PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT)


# ================================================================
# Slide 16: AppPack Add-on
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
bar(s, 0.6, 0.4, 0.08, 0.5, ACCENT)
txt(s, 0.8, 0.4, 12, 0.7, "AppPack 增值选配", 36, WHITE, True)

packs = [
    ("能源行业包", "¥15-30万/年", "电力调度\n设备监控\n能耗优化\n安全预警", ACCENT, DARK_BLUE),
    ("科研行业包", "¥10-25万/年", "实验管理\n数据分析\n文献智检\n科研协作", GREEN, DARK_GREEN),
    ("制造行业包", "¥15-30万/年", "质量管控\n供应链决策\n生产排程\n设备预测", ORANGE, DARK_RED),
    ("定制 Pack", "按需报价", "客户专属场景\n深度定制开发\n知识沉淀\n持续迭代", PURPLE, DARK_PURPLE),
]
for i, (name, price, desc, clr, bg_clr) in enumerate(packs):
    x = 0.6 + i * 3.15
    card(s, x, 1.3, 2.9, 5.0, bg_clr)
    bar(s, x, 1.3, 2.9, 0.06, clr)
    txt(s, x + 0.2, 1.6, 2.5, 0.5, name, 22, clr, True)
    txt(s, x + 0.2, 2.2, 2.5, 0.5, price, 20, WHITE, True)
    txt(s, x + 0.2, 2.9, 2.5, 3.0, desc, 16, LIGHT)

card(s, 0.6, 6.5, 12.1, 0.7, CARD)
txt(s, 1.0, 6.55, 11.5, 0.5, "AppPack 跨部署模式通用 — 一体机和集群使用同一套 Pack，迁移无需重新开发", 16, ACCENT2)


# ================================================================
# Slide 17: Migration Path
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
bar(s, 0.6, 0.4, 0.08, 0.5, ACCENT)
txt(s, 0.8, 0.4, 12, 0.7, "架构迁移路径", 36, WHITE, True)

# Path nodes
nodes = [
    (0.8, 2.0, "Box Lite\n¥10万", ACCENT, DARK_BLUE),
    (3.8, 2.0, "Box Pro\n¥30万", ORANGE, DARK_RED),
    (3.8, 4.2, "Box Gov\n¥45万", PURPLE, DARK_PURPLE),
    (7.3, 2.0, "Cluster\nStandard\n¥100-200万", GREEN, DARK_GREEN),
    (7.3, 4.2, "Cluster\nGov\n¥200-400万", RED, DARK_RED),
    (10.8, 2.0, "大规模集群\n5-10 节点", ACCENT2, CARD),
]
for x, y, label, clr, bg_clr in nodes:
    card(s, x, y, 2.5, 1.5, bg_clr)
    bar(s, x, y, 2.5, 0.05, clr)
    txt(s, x + 0.1, y + 0.2, 2.3, 1.2, label, 15, WHITE, True, PP_ALIGN.CENTER)

# Arrow labels
txt(s, 2.0, 2.4, 1.5, 0.4, "加装 GPU →", 13, ACCENT2)
txt(s, 5.5, 2.4, 1.5, 0.4, "数据迁移 →", 13, ACCENT2)
txt(s, 9.5, 2.4, 1.2, 0.4, "加节点 →", 13, ACCENT2)
txt(s, 2.0, 4.5, 1.5, 0.4, "信创替换 ↓", 13, PURPLE)
txt(s, 5.5, 4.5, 1.5, 0.4, "信创替换 →", 13, PURPLE)

# Migration guarantees
card(s, 0.6, 6.0, 12.1, 1.2, CARD)
txt(s, 1.0, 6.1, 11, 0.4, "迁移保障", 20, ACCENT, True)
guarantees = "dts-cli migrate 全量+增量同步  |  CRD 配置导出/导入  |  AppPack 跨模式通用  |  迁移前自动快照  |  失败一键回滚"
txt(s, 1.0, 6.5, 11.5, 0.5, guarantees, 15, LIGHT)


# ================================================================
# Slide 18: Compatibility Matrix
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
bar(s, 0.6, 0.4, 0.08, 0.5, ACCENT)
txt(s, 0.8, 0.4, 12, 0.7, "架构兼容性矩阵", 36, WHITE, True)

# Headers
m_headers = ["特性", "Lite", "Pro", "Gov", "Cluster Std", "Cluster Gov"]
m_widths = [2.5, 1.6, 1.6, 1.6, 2.2, 2.2]
m_x = [0.6]
for w in m_widths[:-1]:
    m_x.append(m_x[-1] + w + 0.1)

for i, (h, x, w) in enumerate(zip(m_headers, m_x, m_widths)):
    card(s, x, 1.3, w, 0.6, DARK_BLUE)
    txt(s, x + 0.1, 1.35, w - 0.2, 0.5, h, 16, ACCENT if i == 0 else WHITE, True, PP_ALIGN.CENTER)

m_rows = [
    ("25 微服务全栈", "Y", "Y", "Y", "Y", "Y"),
    ("本地 LLM", "CPU", "Y", "Y", "Y", "Y"),
    ("HA 高可用", "—", "—", "—", "Y", "Y"),
    ("中间件 HA", "—", "—", "—", "Y", "Y"),
    ("信创认证", "—", "—", "Y", "—", "Y"),
    ("离线部署", "Y", "Y", "Y", "Y", "Y"),
    ("AppPack", "Y", "Y", "Y", "Y", "Y"),
]
for r, row in enumerate(m_rows):
    y = 2.1 + r * 0.6
    if r % 2 == 0:
        bar(s, 0.6, y, 11.7, 0.55, RGBColor(0x1E, 0x1E, 0x32))
    for c, (val, x, w) in enumerate(zip(row, m_x, m_widths)):
        if c == 0:
            txt(s, x + 0.1, y + 0.08, w - 0.2, 0.4, val, 15, GRAY, True)
        else:
            clr = GREEN if val == "Y" else (ORANGE if val == "CPU" else RGBColor(0x55, 0x55, 0x66))
            txt(s, x + 0.1, y + 0.08, w - 0.2, 0.4, val, 16, clr, True, PP_ALIGN.CENTER)


# ================================================================
# Slide 19: Delivery Timeline
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
bar(s, 0.6, 0.4, 0.08, 0.5, ACCENT)
txt(s, 0.8, 0.4, 12, 0.7, "交付流程", 36, WHITE, True)

# Box delivery
card(s, 0.6, 1.3, 5.8, 5.5, DARK_BLUE)
txt(s, 1.0, 1.5, 5, 0.5, "一体机交付 — 5 天", 24, ACCENT, True)
bar(s, 1.0, 2.1, 4.5, 0.04, ACCENT)

box_days = [
    ("Day 0", "硬件出厂预装 DTS + 质检", ACCENT),
    ("Day 1", "到场通电 → 初始化向导 (30min)\n接入 1-2 个核心数据源", WHITE),
    ("Day 2-3", "配置行业 AppPack\n导入业务知识", WHITE),
    ("Day 4-5", "用户培训 + 验收", GREEN),
]
for i, (day, desc, clr) in enumerate(box_days):
    y = 2.4 + i * 1.0
    circle(s, 1.2, y, 0.15, clr)
    txt(s, 1.6, y - 0.1, 1.2, 0.4, day, 16, clr, True)
    txt(s, 2.8, y - 0.1, 3.2, 0.8, desc, 14, LIGHT)

# Cluster delivery
card(s, 6.9, 1.3, 5.8, 5.5, DARK_GREEN)
txt(s, 7.3, 1.5, 5, 0.5, "集群交付 — 10 天", 24, GREEN, True)
bar(s, 7.3, 2.1, 4.5, 0.04, GREEN)

cluster_days = [
    ("Day 0", "硬件到场上架", GREEN),
    ("Day 1", "K8s 集群初始化 (dts-cli install)", WHITE),
    ("Day 2", "DTS 全栈部署 + HA 验证", WHITE),
    ("Day 3-5", "数据源接入 + AppPack 配置", WHITE),
    ("Day 6-8", "压力测试 + HA 故障演练", ORANGE),
    ("Day 9-10", "用户培训 + 验收", GREEN),
]
for i, (day, desc, clr) in enumerate(cluster_days):
    y = 2.4 + i * 0.78
    circle(s, 7.5, y, 0.15, clr)
    txt(s, 7.9, y - 0.1, 1.2, 0.4, day, 14, clr, True)
    txt(s, 9.1, y - 0.1, 3.2, 0.6, desc, 13, LIGHT)


# ================================================================
# Slide 20: Pricing Integration
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
bar(s, 0.6, 0.4, 0.08, 0.5, ACCENT)
txt(s, 0.8, 0.4, 12, 0.7, "与现有定价体系的关系", 36, WHITE, True)

# Left: existing model
card(s, 0.6, 1.4, 5.5, 4.0, DARK_BLUE)
txt(s, 1.0, 1.6, 5, 0.5, "现有三层定价（纯软件）", 22, ACCENT, True)
bar(s, 1.0, 2.2, 4.5, 0.04, ACCENT)

layers_left = [
    ("L1: 平台基础费", "按部署规模，年费 30-300 万", ACCENT),
    ("L2: AppPack 费", "按行业场景，独立定价", GREEN),
    ("L3: 知识分润", "客户知识资产，平台抽成 30%", ORANGE),
]
for i, (title, desc, clr) in enumerate(layers_left):
    y = 2.5 + i * 0.85
    bar(s, 1.0, y, 0.08, 0.6, clr)
    txt(s, 1.3, y, 2.5, 0.4, title, 17, clr, True)
    txt(s, 1.3, y + 0.35, 4.0, 0.4, desc, 14, LIGHT)

# Arrow
txt(s, 6.3, 3.0, 0.6, 0.5, "→", 48, ACCENT2, True, PP_ALIGN.CENTER)

# Right: hardware bundled
card(s, 7.0, 1.4, 5.5, 4.0, DARK_GREEN)
txt(s, 7.4, 1.6, 5, 0.5, "硬件捆绑定价（本方案）", 22, GREEN, True)
bar(s, 7.4, 2.2, 4.5, 0.04, GREEN)

layers_right = [
    ("L1 → 整机报价", "平台费内含在硬件捆绑价中", GREEN),
    ("L2 → 可选增值", "AppPack 独立报价，按需选配", ORANGE),
    ("L3 → 不变", "知识资产分润机制不变", ACCENT2),
]
for i, (title, desc, clr) in enumerate(layers_right):
    y = 2.5 + i * 0.85
    bar(s, 7.4, y, 0.08, 0.6, clr)
    txt(s, 7.7, y, 2.5, 0.4, title, 17, clr, True)
    txt(s, 7.7, y + 0.35, 4.0, 0.4, desc, 14, LIGHT)

# Summary
card(s, 0.6, 5.8, 12.1, 1.2, CARD)
txt(s, 1.0, 5.9, 5.5, 0.5, "纯软件客户:", 18, GRAY, True)
txt(s, 3.5, 5.9, 8, 0.5, "L1 + L2 + L3", 18, ACCENT)
txt(s, 1.0, 6.4, 5.5, 0.5, "硬件捆绑客户:", 18, GRAY, True)
txt(s, 3.5, 6.4, 8, 0.5, "整机价 (含 L1) + L2 (可选) + L3", 18, GREEN)


# ================================================================
# Slide 21: Hardware Partnership (TBD)
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
bar(s, 0.6, 0.4, 0.08, 0.5, ACCENT)
txt(s, 0.8, 0.4, 12, 0.7, "硬件合作模式（待定）", 36, WHITE, True)

partners = [
    ("OEM 贴牌", "品牌统一，利润高\n库存风险大", "规模化阶段", ACCENT, DARK_BLUE),
    ("渠道合作", "轻资产，借力渠道\n利润分薄", "快速扩张", GREEN, DARK_GREEN),
    ("BOM 认证", "最灵活，零库存\n交付一致性难控", "早期验证", ORANGE, DARK_RED),
    ("混合模式", "按需灵活组合\n管理较复杂", "成熟阶段", PURPLE, DARK_PURPLE),
]
for i, (name, desc, phase, clr, bg_clr) in enumerate(partners):
    x = 0.6 + i * 3.15
    card(s, x, 1.4, 2.9, 3.5, bg_clr)
    bar(s, x, 1.4, 2.9, 0.06, clr)
    txt(s, x + 0.2, 1.7, 2.5, 0.5, name, 22, clr, True)
    txt(s, x + 0.2, 2.4, 2.5, 1.5, desc, 15, LIGHT)
    txt(s, x + 0.2, 3.8, 2.5, 0.4, phase, 14, GRAY)

card(s, 0.6, 5.3, 12.1, 1.5, DARK_BLUE)
txt(s, 1.0, 5.5, 11, 0.5, "当前策略", 22, ACCENT, True)
txt(s, 1.0, 6.0, 11, 0.5, "先以 BOM 认证清单模式启动，积累交付案例后再考虑 OEM / 渠道合作", 18, LIGHT)


# ================================================================
# Slide 22: Summary
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
bar(s, 0.6, 0.4, 0.08, 0.5, ACCENT)
txt(s, 0.8, 0.4, 12, 0.7, "方案总结", 36, WHITE, True)

summary_items = [
    ("5 个 SKU", "覆盖 10 万 → 400 万全价格带", ACCENT),
    ("2 种形态", "一体机（塔式/机架）+ 三节点集群", GREEN),
    ("3 种报价", "买断 / 订阅 / 整机方案价，灵活投标", ORANGE),
    ("信创全覆盖", "国产 CPU + GPU + OS，信创目录认证", PURPLE),
    ("平滑迁移", "Lite → Pro → Cluster，数据零丢失", ACCENT2),
    ("开箱即用", "一体机 Day 1 部署 < 1 小时", YELLOW),
]
for i, (title, desc, clr) in enumerate(summary_items):
    row = i // 2
    col = i % 2
    x = 0.6 + col * 6.3
    y = 1.4 + row * 1.7
    card(s, x, y, 5.9, 1.4, CARD)
    bar(s, x, y, 0.08, 1.4, clr)
    txt(s, x + 0.3, y + 0.15, 5.3, 0.5, title, 22, clr, True)
    txt(s, x + 0.3, y + 0.65, 5.3, 0.5, desc, 16, LIGHT)

txt(s, 0.8, 6.5, 12, 0.5, "DTS — 让每一个决策都有数字孪生  |  软硬一体，开箱即用", 20, GRAY, False, PP_ALIGN.CENTER)


# ================================================================
# Save
# ================================================================
out = os.path.join(os.path.dirname(__file__), "DTS-部署与定价方案-V1.0.0.pptx")
prs.save(out)
print(f"Saved: {out}")
