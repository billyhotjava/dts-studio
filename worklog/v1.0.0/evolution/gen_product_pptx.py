#!/usr/bin/env python3
"""Generate DTS v3.0 Product Specification PPT — 产品功能说明书."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Colors ──
BG = RGBColor(0x1A, 0x1A, 0x2E)
CARD = RGBColor(0x22, 0x22, 0x3A)
ACCENT = RGBColor(0x00, 0x96, 0xD6)
ACCENT2 = RGBColor(0x00, 0xC9, 0xA7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xAA, 0xAA, 0xBB)
LIGHT = RGBColor(0xE0, 0xE0, 0xE8)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
RED = RGBColor(0xE8, 0x4D, 0x4D)
GREEN = RGBColor(0x28, 0xC7, 0x6F)
PURPLE = RGBColor(0xBB, 0x86, 0xFC)
DARK_BLUE = RGBColor(0x0A, 0x3D, 0x5C)
DARK_RED = RGBColor(0x3A, 0x1A, 0x1A)
DARK_GREEN = RGBColor(0x0A, 0x3D, 0x2A)
DARK_PURPLE = RGBColor(0x2A, 0x1A, 0x3A)
YELLOW = RGBColor(0xFF, 0xD7, 0x00)


def bg(slide):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = BG

def txt(slide, l, t, w, h, text, sz=18, clr=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = clr; p.font.bold = bold
    p.font.name = "Microsoft YaHei"; p.alignment = align
    return tb

def card(slide, l, t, w, h, clr=CARD):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = clr; s.line.fill.background(); s.shadow.inherit = False
    return s

def bullets(slide, l, t, w, h, items, sz=16, clr=LIGHT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(sz); p.font.color.rgb = clr
        p.font.name = "Microsoft YaHei"; p.space_after = Pt(8)
    return tb

def bar(slide, l, t, w, h, clr):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = clr; s.line.fill.background()
    return s

def circle(slide, l, t, d, clr):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(d), Inches(d))
    s.fill.solid(); s.fill.fore_color.rgb = clr; s.line.fill.background()
    return s


# ════════════════════════════════════════════════════════
# Slide 1: Cover
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
bar(s, 1, 2.3, 1.5, 0.06, ACCENT)
txt(s, 1, 2.6, 11, 1.2, "Decision Twins System", 48, WHITE, True)
txt(s, 1, 3.7, 11, 0.8, "产品功能说明书 v3.0", 28, ACCENT)
txt(s, 1, 4.6, 8, 0.6, "AI Decision Operating System", 20, GRAY)
txt(s, 1, 5.3, 8, 0.5, "Ontology + Intent + Agent = Decision Twins", 16, ACCENT2)
txt(s, 1, 6.2, 8, 0.4, "2026  |  DRAFT", 14, GRAY)


# ════════════════════════════════════════════════════════
# Slide 2: Product Overview
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.8, 0.4, 12, 0.7, "产品定位", 36, WHITE, True)

card(s, 0.8, 1.3, 11.7, 1.8, DARK_BLUE)
txt(s, 1.2, 1.5, 11, 0.5,
    "DTS 是 AI 驱动的企业决策操作系统", 24, WHITE, True)
txt(s, 1.2, 2.0, 11, 0.9,
    "不是数据中台（不卖存储/计算）| 不是 BI（不只看数据）| 不是 SaaS（不按人头卖工具）\n"
    "而是把数据资产、业务规则、行业隐性知识编织成一个能替决策者「想」的系统", 17, LIGHT)

txt(s, 0.8, 3.4, 12, 0.5, "核心价值公式", 20, WHITE, True)
card(s, 0.8, 3.9, 11.7, 1.2, DARK_PURPLE)
txt(s, 1.2, 4.05, 11, 0.5,
    "1 位专家的创造力  ×  AI 提取能力  ×  系统化规模  =  500 倍产出放大", 22, WHITE, True, PP_ALIGN.CENTER)
txt(s, 1.2, 4.55, 11, 0.4,
    "用户一句话完成：理解意图 → 查数据 → 分析 → 图表 → 报告 → 建议 → 执行", 16, LIGHT, False, PP_ALIGN.CENTER)

# Product forms
txt(s, 0.8, 5.4, 12, 0.5, "产品形态", 20, WHITE, True)
forms = [
    ("三个前端应用", "Pilot 决策门户\nCortex 专家工作台\nAdmin 管理控制台", ACCENT),
    ("开放 API 体系", "REST (外部)\ngRPC (内部)\nSSE (流式推送)", ACCENT2),
    ("AppPack 生态", "标准化行业能力包\nOntology + Skill + KnowHow\n+ 配置 + UI 全栈插件", GREEN),
]
for i, (title, desc, clr) in enumerate(forms):
    x = 0.8 + i * 4.1
    card(s, x, 5.9, 3.8, 1.4)
    bar(s, x, 5.9, 3.8, 0.06, clr)
    txt(s, x + 0.2, 6.05, 3.4, 0.35, title, 17, clr, True)
    txt(s, x + 0.2, 6.4, 3.4, 0.8, desc, 14, LIGHT)


# ════════════════════════════════════════════════════════
# Slide 3: Core Formula — Ontology + Intent + Agent
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.8, 0.4, 12, 0.7, "核心三元组：Ontology + Intent + Agent", 32, WHITE, True)

components = [
    ("Ontology", "认知基座", "物理世界到数字世界的映射标准\n不只是数据字典，而是包含\n实体、关系、指标、规则、操作的\n完整认知框架", ACCENT),
    ("Intent", "意图理解", "自然语言到结构化决策意图\n支持多目标分解\n约束识别\n风险评估", ACCENT2),
    ("Agent", "自主执行", "基于 LangGraph 的多步推理运行时\n自主规划执行路径\n调用技能\n生成决策解释", GREEN),
]
for i, (name, subtitle, desc, clr) in enumerate(components):
    x = 0.8 + i * 4.1
    card(s, x, 1.3, 3.8, 3.5)
    bar(s, x, 1.3, 3.8, 0.08, clr)
    txt(s, x + 0.3, 1.55, 3.2, 0.5, name, 28, clr, True)
    txt(s, x + 0.3, 2.05, 3.2, 0.4, subtitle, 18, WHITE)
    txt(s, x + 0.3, 2.55, 3.2, 2.0, desc, 15, LIGHT)

# Connection
card(s, 0.8, 5.2, 11.7, 1.8, DARK_BLUE)
txt(s, 1.2, 5.35, 11, 0.5,
    "三者通过 DAP (DTS Agent Protocol) 协议连接，形成完整的决策闭环", 20, WHITE, True)
txt(s, 1.2, 5.85, 11, 1.0,
    "用户说「帮我看看A3产线今天的OEE，如果低于80%就安排维保」\n"
    "→ Intent 解析为多目标 → Agent 规划执行步骤 → 查 Ontology 找指标定义 → 调 Skill 查数据\n"
    "→ 评估条件 → 触发维保审批 → 全程审计 → 流式返回结果", 15, LIGHT)


# ════════════════════════════════════════════════════════
# Slide 4: Target User Roles
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.8, 0.4, 12, 0.7, "目标用户角色", 36, WHITE, True)

roles = [
    ("决策者", "高管、部门负责人", "用自然语言获取分析结果\n和行动建议", "Pilot", ACCENT),
    ("行业专家", "分析师、工程师", "建模、探索数据\n编排技能、管理知识", "Cortex", ACCENT2),
    ("IT 管理员", "系统管理员、运维", "管理权限、配置数据源\n监控系统", "Admin", GREEN),
    ("开发者", "ISV、集成商", "开发 AppPack\n调用 API、扩展能力", "ADK + API", ORANGE),
]
for i, (role, title, need, app, clr) in enumerate(roles):
    x = 0.8 + i * 3.1
    card(s, x, 1.2, 2.8, 4.5)
    bar(s, x, 1.2, 2.8, 0.08, clr)
    circle(s, x + 0.9, 1.5, 1.0, clr)
    txt(s, x + 0.9, 1.65, 1.0, 0.7, role[:1], 32, WHITE, True, PP_ALIGN.CENTER)
    txt(s, x + 0.2, 2.7, 2.4, 0.4, role, 20, clr, True)
    txt(s, x + 0.2, 3.1, 2.4, 0.3, title, 14, GRAY)
    txt(s, x + 0.2, 3.5, 2.4, 1.0, need, 15, LIGHT)
    txt(s, x + 0.2, 4.6, 2.4, 0.4, f"使用端: {app}", 14, clr, True)


# ════════════════════════════════════════════════════════
# Slide 5: Three-Pillar Architecture
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.8, 0.4, 12, 0.7, "三支柱架构", 36, WHITE, True)

pillars = [
    ("dts-studio", "设计中心", "Ontology 可视化设计\nSkill 编排\nKnowHow 管理\nPack 开发 IDE", ACCENT),
    ("dts-stack", "Agent 能力中心", "25 个微服务\n57 个 AI 技能\nDAP 协议运行时\n五条铁律保障", ACCENT2),
    ("app-stack", "客户 Agent 中心", "行业 AppPack\n客户定制 Pack\nKnowHow 资产\nISV 开发的 Pack", GREEN),
]
for i, (name, subtitle, desc, clr) in enumerate(pillars):
    x = 0.8 + i * 4.1
    card(s, x, 1.2, 3.8, 3.0)
    bar(s, x, 1.2, 3.8, 0.08, clr)
    txt(s, x + 0.3, 1.45, 3.2, 0.5, name, 24, clr, True)
    txt(s, x + 0.3, 1.95, 3.2, 0.35, subtitle, 16, WHITE)
    txt(s, x + 0.3, 2.4, 3.2, 1.5, desc, 15, LIGHT)

# Three-language collaboration
txt(s, 0.8, 4.5, 12, 0.5, "三语言协作", 22, WHITE, True)
langs = [
    ("Python", "10 个 AI 核心服务", "Agent/Intent/RAG/Eval\n全部可降级", ACCENT),
    ("Java", "10 个业务平台服务", "CRUD/RBAC/安全/存储\n独立可运行", ACCENT2),
    ("Go", "2 个基础设施服务", "K8s Operator + CLI\n部署运维自动化", GREEN),
    ("React", "3 个前端应用", "Pilot + Cortex + Admin\nReact 19 + Vite + shadcn", ORANGE),
]
for i, (lang, count, desc, clr) in enumerate(langs):
    x = 0.8 + i * 3.1
    card(s, x, 5.0, 2.8, 2.2)
    txt(s, x + 0.2, 5.1, 2.4, 0.35, lang, 20, clr, True)
    txt(s, x + 0.2, 5.45, 2.4, 0.3, count, 14, WHITE)
    txt(s, x + 0.2, 5.8, 2.4, 1.0, desc, 14, LIGHT)


# ════════════════════════════════════════════════════════
# Slide 6: DAP Protocol
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.8, 0.4, 12, 0.7, "DAP 协议 — 产品核心护城河", 36, WHITE, True)
txt(s, 0.8, 1.0, 12, 0.4, "DTS Agent Protocol: 定义 AI Agent 如何理解物理世界并执行决策", 16, GRAY)

layers = [
    ("Layer 5", "Audit Protocol", "决策可追溯性保障 (CloudEvents, append-only)", PURPLE),
    ("Layer 4", "Security Protocol", "身份与权限在 Agent 间传递 (JWT + 零信任)", RED),
    ("Layer 3", "Intent Protocol", "自然语言 → 结构化意图 → 执行计划", ORANGE),
    ("Layer 2", "Skill Protocol", "AI 技能的定义、注册、调用、编排 (gRPC streaming)", GREEN),
    ("Layer 1", "Ontology Protocol", "物理世界 → 数字映射的标准格式 (YAML 声明式)", ACCENT2),
    ("Layer 0", "Transport", "gRPC (control) + Kafka (data) + SSE (streaming)", ACCENT),
]
for i, (layer, name, desc, clr) in enumerate(layers):
    y = 1.5 + i * 0.85
    card(s, 0.8, y, 11.7, 0.75)
    bar(s, 0.8, y, 0.12, 0.75, clr)
    txt(s, 1.2, y + 0.08, 1.5, 0.35, layer, 16, clr, True)
    txt(s, 2.7, y + 0.08, 3.0, 0.35, name, 16, WHITE, True)
    txt(s, 5.8, y + 0.08, 6.5, 0.6, desc, 14, LIGHT)

card(s, 0.8, 6.7, 11.7, 0.6, DARK_BLUE)
txt(s, 1.2, 6.75, 11, 0.4,
    "协议版本: dap/v1 (当前) → v2 (多模态) → v3 (跨组织 Agent 协作)  |  向后兼容", 15, LIGHT)


# ════════════════════════════════════════════════════════
# Slide 7: Core Data Flow
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.8, 0.4, 12, 0.7, "核心数据流", 36, WHITE, True)

steps = [
    ("用户输入", "自然语言", ACCENT),
    ("Intent 解析", "NL → 结构化\ngoals + constraints", ACCENT2),
    ("Agent 规划", "Intent → Plan\n多步骤 DAG", GREEN),
    ("Skill 执行", "gRPC streaming\n逐步执行", ORANGE),
    ("结果输出", "富卡片/报告\n图表/Action", PURPLE),
]
for i, (step, desc, clr) in enumerate(steps):
    x = 0.8 + i * 2.46
    card(s, x, 1.2, 2.2, 2.0)
    bar(s, x, 1.2, 2.2, 0.06, clr)
    txt(s, x + 0.15, 1.4, 1.9, 0.35, step, 16, clr, True)
    txt(s, x + 0.15, 1.8, 1.9, 1.0, desc, 14, LIGHT)
    if i < 4:
        txt(s, x + 2.15, 1.7, 0.3, 0.5, "→", 24, GRAY)

# Risk branching
txt(s, 0.8, 3.5, 12, 0.5, "风险分级与人工介入", 20, WHITE, True)
risk_levels = [
    ("低风险", "自动执行\nSSE 推送进度", GREEN),
    ("中风险", "暂停 → HITL 审批\n审批通过 → 继续", ORANGE),
    ("高风险", "强制暂停 → 多级审批\n展示 AI 推理链路", RED),
]
for i, (level, desc, clr) in enumerate(risk_levels):
    x = 0.8 + i * 4.1
    card(s, x, 4.0, 3.8, 1.5)
    bar(s, x, 4.0, 3.8, 0.06, clr)
    txt(s, x + 0.3, 4.15, 3.2, 0.35, level, 18, clr, True)
    txt(s, x + 0.3, 4.55, 3.2, 0.8, desc, 15, LIGHT)

# Audit
card(s, 0.8, 5.8, 11.7, 1.4, DARK_BLUE)
txt(s, 1.2, 5.95, 11, 0.4, "全链路审计", 20, WHITE, True)
txt(s, 1.2, 6.35, 11, 0.7,
    "全部操作 → Kafka → dts-audit-log (append-only, 不可篡改) + dts-observability (AI Trace 树, 成本统计)\n"
    "人和 AI 走同一个 Capability API，审计天然统一，无需额外实现", 15, LIGHT)


# ════════════════════════════════════════════════════════
# Slide 8: Five Iron Laws in Product
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.8, 0.4, 12, 0.7, "五条铁律在产品中的体现", 36, WHITE, True)

laws = [
    ("1", "人工随时可接管", "每个 AI 能力都有等价的人工操作路径，AI 不可用时平台 100% 可用", ACCENT),
    ("2", "核心安全不可绕过", "所有请求必经 dts-gateway，AppPack 无特权，Skill 调用携带 RequestContext", ACCENT2),
    ("3", "数据安全是基因", "每个属性标记 classification，AI 读取自动过 dts-data-security，SQL 自动注入行级过滤", GREEN),
    ("4", "全操作可追溯", "人和 AI 走同一个 Capability API，SkillResponse 携带审计条目，Kafka append-only", ORANGE),
    ("5", "能力先于界面", "每个服务先有完整 API + 测试，AI Tool 是 API 适配器，UI 是加速层不是替代层", PURPLE),
]
for i, (num, title, desc, clr) in enumerate(laws):
    y = 1.2 + i * 1.15
    card(s, 0.8, y, 11.7, 1.0)
    circle(s, 1.1, y + 0.15, 0.7, clr)
    txt(s, 1.1, y + 0.2, 0.7, 0.6, num, 26, WHITE, True, PP_ALIGN.CENTER)
    txt(s, 2.1, y + 0.1, 3.0, 0.4, title, 20, WHITE, True)
    txt(s, 2.1, y + 0.5, 9.5, 0.5, desc, 14, LIGHT)


# ════════════════════════════════════════════════════════
# Slide 9: Pilot — Decision Portal
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.8, 0.4, 12, 0.7, "Pilot — 决策门户", 36, WHITE, True)
txt(s, 0.8, 1.0, 8, 0.4, "面向管理者和业务人员的自然语言交互式决策入口", 16, GRAY)

features = [
    ("自然语言交互", "AI 对话主界面 + CopilotSidebar\n流式响应 (SSE)\n对话内嵌富卡片\n追问式增量修改", ACCENT),
    ("决策仪表盘", "KPI 概览卡片\nAI 图表 Pin-to-Canvas\n保存为报告/看板\n角色自适应菜单", ACCENT2),
    ("审批工作流", "高风险操作推送审批\n展示 AI 推理链路\n多级审批流\n超时自动处理", GREEN),
    ("通知与推送", "定时报告推送\n异常告警推送\n多渠道投递\n移动端适配", ORANGE),
]
for i, (title, desc, clr) in enumerate(features):
    x = 0.8 + i * 3.1
    card(s, x, 1.5, 2.8, 3.5)
    bar(s, x, 1.5, 2.8, 0.06, clr)
    txt(s, x + 0.2, 1.65, 2.4, 0.4, title, 18, clr, True)
    txt(s, x + 0.2, 2.1, 2.4, 2.5, desc, 15, LIGHT)

card(s, 0.8, 5.3, 11.7, 2.0, DARK_BLUE)
txt(s, 1.2, 5.45, 11, 0.4, "核心场景演示", 20, WHITE, True)
txt(s, 1.2, 5.85, 11, 0.4,
    "用户说：「统计华东区半年销售，出报告」", 18, ACCENT, True)
txt(s, 1.2, 6.3, 11, 0.8,
    "理解意图 → 自动生成查询 → 从数据库获取 → 月度趋势+同比分析 → 折线图+柱状图 → 标题+结论+建议 → 一键导出 PDF", 16, LIGHT)


# ════════════════════════════════════════════════════════
# Slide 10: Cortex — Expert Workbench
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.8, 0.4, 12, 0.7, "Cortex — 专家工作台", 36, WHITE, True)
txt(s, 0.8, 1.0, 10, 0.4, "面向行业专家的深度工作环境：数据探索 → Ontology 建模 → 技能编排 → 知识管理", 16, GRAY)

features = [
    ("Ontology 可视化", "React Flow 图可视化\n拖拽式定义\n语义搜索 (pgvector)\n数据血缘追踪", ACCENT),
    ("KnowHow 管理", "按金字塔层级展示\n专家校正 AI 草稿\n沙箱试运行\n碎片管理和聚类", ACCENT2),
    ("技能编排工作台", "DAG 拖拽编排\nSkill 测试运行\nPersona 管理\n评测对比分析", GREEN),
    ("数据探索分析", "SQL 编辑器\n数据 Profiling\n质量看板\nNL2SQL 对话查询", ORANGE),
]
for i, (title, desc, clr) in enumerate(features):
    x = 0.8 + i * 3.1
    card(s, x, 1.5, 2.8, 3.5)
    bar(s, x, 1.5, 2.8, 0.06, clr)
    txt(s, x + 0.2, 1.65, 2.4, 0.4, title, 18, clr, True)
    txt(s, x + 0.2, 2.1, 2.4, 2.5, desc, 15, LIGHT)

txt(s, 0.8, 5.3, 12, 0.4,
    "Cortex 的核心价值：将个人的行业经验转化为可复用的组织资产", 18, ACCENT)


# ════════════════════════════════════════════════════════
# Slide 11: Admin — Management Console
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.8, 0.4, 12, 0.7, "Admin — 管理控制台", 36, WHITE, True)
txt(s, 0.8, 1.0, 10, 0.4, "面向 IT 管理员的平台运维和配置中心", 16, GRAY)

features = [
    ("用户/角色/权限", "RBAC + Keycloak 同步\n多租户隔离\n功能权限与数据权限分离\n角色自适应菜单", ACCENT),
    ("数据源配置", "JDBC/File/API/OPC-UA\n一键连接测试\n元数据自动发现\nCDC 配置", ACCENT2),
    ("服务监控", "25 服务健康看板\nGrafana Tempo 追踪\nPrometheus 指标\nLoki 日志", GREEN),
    ("AppPack 管理", "应用市场浏览/安装\n版本管理/回滚\n运行状态监控\n权限审批", ORANGE),
]
for i, (title, desc, clr) in enumerate(features):
    x = 0.8 + i * 3.1
    card(s, x, 1.5, 2.8, 3.2)
    bar(s, x, 1.5, 2.8, 0.06, clr)
    txt(s, x + 0.2, 1.65, 2.4, 0.4, title, 18, clr, True)
    txt(s, x + 0.2, 2.1, 2.4, 2.3, desc, 15, LIGHT)

# Additional features
extras = [
    ("审计日志查看", "全操作日志 + 多维筛选 + 决策链路回溯 + 合规报告导出", PURPLE),
    ("系统配置", "LLM 端点/Token 限额/超时策略 + 安全策略 + 调度管理 + 备份管理", GRAY),
]
for i, (title, desc, clr) in enumerate(extras):
    y = 5.0 + i * 0.85
    card(s, 0.8, y, 11.7, 0.75)
    txt(s, 1.2, y + 0.1, 2.5, 0.35, title, 16, clr, True)
    txt(s, 3.8, y + 0.1, 8.5, 0.5, desc, 14, LIGHT)


# ════════════════════════════════════════════════════════
# Slide 12: Ontology Management
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.8, 0.4, 12, 0.7, "核心模块：Ontology 管理", 36, WHITE, True)
txt(s, 0.8, 1.0, 12, 0.4,
    "AI Agent 理解客户物理世界的认知框架 — 不是数据字典，是完整数字孪生模型", 16, GRAY)

ont_features = [
    ("ObjectType", "声明式定义业务对象\n属性/类型/分级/来源", ACCENT),
    ("Relationship", "React Flow 图可视化\n关系和基数管理", ACCENT2),
    ("Metric", "公式编辑+阈值告警\nOEE = A × P × Q", GREEN),
    ("Action", "可执行操作定义\n风险等级 → HITL", ORANGE),
]
for i, (name, desc, clr) in enumerate(ont_features):
    x = 0.8 + i * 3.1
    card(s, x, 1.5, 2.8, 1.8)
    bar(s, x, 1.5, 2.8, 0.06, clr)
    txt(s, x + 0.2, 1.65, 2.4, 0.35, name, 18, clr, True)
    txt(s, x + 0.2, 2.05, 2.4, 1.0, desc, 15, LIGHT)

# DAP design principles
txt(s, 0.8, 3.6, 12, 0.4, "DAP 设计原则", 20, WHITE, True)
principles = [
    ("声明式而非命令式", "定义「是什么」，Agent 自主决定「怎么做」", ACCENT),
    ("数据安全内置", "每个属性标记 classification，读取时自动过安全检查", GREEN),
    ("来源标记", "标记 source 和 frequency，Agent 知道去哪拿数据", ORANGE),
    ("行业命名空间", "manufacturing/energy/research 隔离，防跨 Pack 冲突", PURPLE),
]
for i, (title, desc, clr) in enumerate(principles):
    y = 4.0 + i * 0.7
    card(s, 0.8, y, 11.7, 0.6)
    bar(s, 0.8, y, 0.1, 0.6, clr)
    txt(s, 1.2, y + 0.05, 3.0, 0.35, title, 16, clr, True)
    txt(s, 4.3, y + 0.05, 8.0, 0.5, desc, 14, LIGHT)


# ════════════════════════════════════════════════════════
# Slide 13: Intent Engine + Agent Runtime
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.8, 0.4, 12, 0.7, "Intent 引擎 + Agent 运行时", 36, WHITE, True)

# Intent Engine
card(s, 0.8, 1.2, 5.8, 3.2, CARD)
bar(s, 0.8, 1.2, 5.8, 0.06, ACCENT2)
txt(s, 1.1, 1.4, 5.2, 0.4, "Intent 引擎 (dts-intent-engine)", 20, ACCENT2, True)
bullets(s, 1.1, 1.85, 5.2, 2.3, [
    "NL → 结构化 Intent (goals + constraints + plan)",
    "多目标意图分解: 一句话中识别多个目标",
    "约束条件识别: max_steps, timeout, approval_level",
    "风险等级评估: low/medium/high 自动分级",
    "Intent 路由: 自动选择 Skill 组合",
    "上下文保持: 多轮对话增量修改",
], 14, LIGHT)

# Agent Runtime
card(s, 7.0, 1.2, 5.8, 3.2, CARD)
bar(s, 7.0, 1.2, 5.8, 0.06, GREEN)
txt(s, 7.3, 1.4, 5.2, 0.4, "Agent 运行时 (dts-agent)", 20, GREEN, True)
bullets(s, 7.3, 1.85, 5.2, 2.3, [
    "LangGraph 状态机驱动的 DAG 执行引擎",
    "Persona 管理: system prompt + Skill 集 + 定时任务",
    "执行计划生成: LLM → TaskPlan (DAG)",
    "标准 gRPC SkillService 接口调用所有 Skill",
    "决策解释: AI Trace 树 → 人类可读推理链",
    "人工接管: 任意步骤可暂停，手动完成后继续",
    "错误恢复: 重试/跳过/降级/终止",
], 14, LIGHT)

# Key insight
card(s, 0.8, 4.7, 11.7, 1.0, DARK_BLUE)
txt(s, 1.2, 4.85, 11, 0.7,
    "Intent 描述「要什么」| Agent 决定「怎么做」| Skill 执行「做什么」— 三者通过 DAP 协议解耦", 18, WHITE, True, PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
# Slide 14: Data Connection & Quality
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.8, 0.4, 12, 0.7, "数据连接与质量", 36, WHITE, True)
txt(s, 0.8, 1.0, 12, 0.4, "不搬数据、就地使用 — JDBC 直连，SQL 下推执行，数据不离开原位", 16, GRAY)

# Data sources
txt(s, 0.8, 1.5, 6, 0.4, "支持的数据源", 20, WHITE, True)
sources = [
    ("关系数据库", "PG/MySQL/Oracle\nSQL Server/ClickHouse", ACCENT),
    ("大数据", "Hive/Spark SQL\n(via JDBC)", ACCENT2),
    ("文件", "CSV/Excel/JSON\nParquet (MinIO)", GREEN),
    ("API/工业", "REST / OPC-UA\nMQTT (IoT)", ORANGE),
]
for i, (cat, items, clr) in enumerate(sources):
    x = 0.8 + i * 3.1
    card(s, x, 1.9, 2.8, 1.3)
    txt(s, x + 0.2, 1.95, 2.4, 0.3, cat, 15, clr, True)
    txt(s, x + 0.2, 2.3, 2.4, 0.7, items, 14, LIGHT)

# Data quality
txt(s, 0.8, 3.5, 6, 0.4, "AI 数据质量检测", 20, WHITE, True)
quality = [
    ("规则检测", "SQL 质量检查\n完整性/一致性/范围", ACCENT),
    ("异常发现", "统计分布对比\n漂移检测 + LLM 解释", ACCENT2),
    ("修复建议", "LLM 提出纠正方案\n规则化去重/标准化", GREEN),
    ("数据 Profiling", "分布/空值/离群\nLLM 解释异常模式", ORANGE),
]
for i, (name, desc, clr) in enumerate(quality):
    x = 0.8 + i * 3.1
    card(s, x, 3.9, 2.8, 1.3)
    txt(s, x + 0.2, 3.95, 2.4, 0.3, name, 15, clr, True)
    txt(s, x + 0.2, 4.3, 2.4, 0.7, desc, 14, LIGHT)


# ════════════════════════════════════════════════════════
# Slide 15: Security Architecture
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.8, 0.4, 12, 0.7, "安全架构", 36, WHITE, True)
txt(s, 0.8, 1.0, 12, 0.4, "数据安全不是功能模块，而是架构基因", 16, GRAY)

security = [
    ("认证", "Keycloak + JWT", "OIDC/SAML/PKI\nJWT 在服务间通过 gRPC metadata 传递\ndts-gateway 强制验证所有外部请求", ACCENT),
    ("授权", "RBAC + ABAC", "功能权限 (菜单/按钮/API)\n数据权限 (分级 × 属性 → allow/mask/deny)\n角色自适应菜单", ACCENT2),
    ("数据安全", "动态脱敏", "行级过滤 (SecuritySqlRewriter)\n列级控制 + 数据水印\nRAG 检索安全过滤", GREEN),
    ("AI 安全", "零信任", "AppPack 无特权\nPrompt 注入防护\nTaskContext 加密存储", ORANGE),
]
for i, (area, method, desc, clr) in enumerate(security):
    y = 1.5 + i * 1.3
    card(s, 0.8, y, 11.7, 1.15)
    bar(s, 0.8, y, 0.12, 1.15, clr)
    txt(s, 1.2, y + 0.05, 1.8, 0.35, area, 18, clr, True)
    txt(s, 3.0, y + 0.05, 2.5, 0.35, method, 16, WHITE, True)
    txt(s, 5.5, y + 0.05, 6.8, 1.0, desc, 14, LIGHT)

card(s, 0.8, 6.8, 11.7, 0.5, DARK_RED)
txt(s, 1.2, 6.85, 11, 0.35,
    "核心保障: 每次数据访问都通过 DataSecurityCheck 验证，不缓存权限决策，AppPack 的 Skill 看到的数据 = 当前用户有权看到的数据", 14, LIGHT)


# ════════════════════════════════════════════════════════
# Slide 16: 57 AI Skills Overview
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.8, 0.4, 12, 0.7, "AI 技能体系 — 57 个标准化技能", 36, WHITE, True)
txt(s, 0.8, 1.0, 12, 0.4,
    "每个 Skill 遵循 DAP Layer 2，gRPC 标准接口: Describe (发现) + Validate (预校验) + Execute (流式执行)", 15, GRAY)

categories = [
    ("平台技能", "17 个", "数据查询/Ontology 建模\n报告生成/Dashboard\n行动执行/导航", ACCENT,
     "data/execute-query\nontology/search\nquery/nl2sql\nreport/generate\naction/create-workflow"),
    ("数据技能", "6 个", "数据接入/质量检测\n异常发现/CDC 同步\n映射建议", ACCENT2,
     "connector/discover-metadata\nquality/run-rules\nquality/detect-drift"),
    ("AI 技能", "9 个", "Agent 推理/Intent 解析\n评测/Bad Case 飞轮\nPersona 切换", GREEN,
     "agent/route-intent\nagent/plan-steps\neval/run-suite\neval/compare-versions"),
    ("行业技能", "18 个", "电力: 负荷预测/缺陷检测\n科研: 质量分析/NCR 归零\n制造: OEE/排程/根因", ORANGE,
     "energy/predict-load\nresearch/analyze-quality\nmfg/predict-maintenance"),
    ("运维技能", "7 个", "健康检查/日志诊断\n扩缩容/备份恢复\nPack 安装/升级", PURPLE,
     "ops/check-health\nops/diagnose-error\nops/install-pack"),
]
for i, (name, count, desc, clr, examples) in enumerate(categories):
    x = 0.8 + i * 2.46
    card(s, x, 1.5, 2.2, 5.0)
    bar(s, x, 1.5, 2.2, 0.06, clr)
    txt(s, x + 0.1, 1.65, 2.0, 0.35, name, 16, clr, True)
    txt(s, x + 0.1, 2.0, 2.0, 0.3, count, 14, WHITE, True)
    txt(s, x + 0.1, 2.35, 2.0, 1.5, desc, 13, LIGHT)
    txt(s, x + 0.1, 4.0, 2.0, 0.3, "示例:", 12, GRAY)
    txt(s, x + 0.1, 4.3, 2.0, 1.8, examples, 11, GRAY)

# Skill types
card(s, 0.8, 6.8, 11.7, 0.5, DARK_BLUE)
txt(s, 1.2, 6.85, 11, 0.35,
    "技能类型: Rule (纯规则) | LLM (大模型推理) | Hybrid (规则+LLM) | Workflow (含人工审批)", 15, LIGHT)


# ════════════════════════════════════════════════════════
# Slide 17: AppPack Ecosystem
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.8, 0.4, 12, 0.7, "AppPack 生态", 36, WHITE, True)
txt(s, 0.8, 1.0, 12, 0.4,
    "半成品 + Connector 策略: 不做完整行业应用，封装决策层精华 (20%)", 16, GRAY)

# 10 capability types
txt(s, 0.8, 1.5, 12, 0.4, "Pack 包含 10 种能力类型", 18, WHITE, True)
cap_types = [
    "ontology", "skills", "connectors", "frontend", "workflows",
    "persona", "metrics", "quality_rules", "actions", "evaluations",
]
for i, cap in enumerate(cap_types):
    row = i // 5
    col = i % 5
    x = 0.8 + col * 2.46
    y = 1.9 + row * 0.55
    card(s, x, y, 2.2, 0.45)
    txt(s, x + 0.1, y + 0.05, 2.0, 0.3, cap, 14, ACCENT, True, PP_ALIGN.CENTER)

# Vertical depth spectrum
txt(s, 0.8, 3.2, 12, 0.4, "垂直深度光谱", 18, WHITE, True)
depths = [
    ("Level 0", "通用数据平台", "传统大数据平台停在这里", GRAY),
    ("Level 1", "行业知识插件", "DTS 核心层: Ontology + Skills + Prompts", GREEN),
    ("Level 2", "垂直应用模块", "未来可选: 在 Pack 里做决策层应用", ACCENT),
    ("Level 3", "独立垂直 SaaS", "不做: 避免变成应用厂商", RED),
]
for i, (level, name, desc, clr) in enumerate(depths):
    y = 3.6 + i * 0.6
    card(s, 0.8, y, 11.7, 0.5)
    txt(s, 1.1, y + 0.05, 1.3, 0.3, level, 14, clr, True)
    txt(s, 2.5, y + 0.05, 2.5, 0.3, name, 14, WHITE, True)
    txt(s, 5.2, y + 0.05, 7.0, 0.35, desc, 13, LIGHT)

# Security
card(s, 0.8, 6.1, 11.7, 1.2, DARK_BLUE)
txt(s, 1.2, 6.2, 11, 0.35, "AppPack 安全保障", 18, WHITE, True)
txt(s, 1.2, 6.55, 11, 0.6,
    "运行时隔离 (GraalVM 沙箱) | 权限沙箱 (DataAccessGate) | 发布签名校验 + 静态扫描\n"
    "网络隔离 (白名单域名) | 与内置 Skill 完全相同的安全链路", 15, LIGHT)


# ════════════════════════════════════════════════════════
# Slide 18: KnowHow Knowledge Management
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.8, 0.4, 12, 0.7, "KnowHow 知识管理", 36, WHITE, True)

# Five-layer pyramid
txt(s, 0.8, 1.1, 6, 0.4, "五层知识金字塔", 20, WHITE, True)
pyramid = [
    ("L4", "世界模型", "Agent 预测操作后果", "TBD", PURPLE),
    ("L3", "隐性模式", "专家没说但一直用的规则", "~60%", ORANGE),
    ("L2", "潜在知识", "专家能说但没人问的逻辑", "~75%", GREEN),
    ("L1", "显性结构", "数据字典、行业标准", "~95%", ACCENT2),
    ("L0", "多模态感知", "视觉+声音+时序+文本", "TBD", ACCENT),
]
for i, (level, name, desc, acc, clr) in enumerate(pyramid):
    y = 1.5 + i * 0.7
    w = 11.7 - i * 0.6
    x = 0.8 + i * 0.3
    card(s, x, y, w, 0.6)
    bar(s, x, y, 0.08, 0.6, clr)
    txt(s, x + 0.2, y + 0.05, 0.6, 0.3, level, 14, clr, True)
    txt(s, x + 0.9, y + 0.05, 2.0, 0.3, name, 14, WHITE, True)
    txt(s, x + 3.0, y + 0.05, 4.0, 0.3, desc, 13, LIGHT)
    txt(s, x + w - 1.5, y + 0.05, 1.3, 0.3, f"准确率: {acc}", 12, GRAY)

# Extraction pipeline
txt(s, 0.8, 5.2, 12, 0.4, "KnowHow 提取四阶段 (Gen-Ba 碎片策略)", 18, WHITE, True)
phases = [
    ("影子观察", "AI 静默观察\n专家日常决策", ACCENT),
    ("碎片外化", "决策瞬间追问\n容许不完整", ACCENT2),
    ("松散组合", "LLM 聚类补全\n生成草稿", GREEN),
    ("协作内化", "与专家并行运行\n评测量化", ORANGE),
]
for i, (name, desc, clr) in enumerate(phases):
    x = 0.8 + i * 3.1
    card(s, x, 5.6, 2.8, 1.3)
    txt(s, x + 0.2, 5.65, 2.4, 0.3, name, 16, clr, True)
    txt(s, x + 0.2, 6.0, 2.4, 0.7, desc, 14, LIGHT)
    if i < 3:
        txt(s, x + 2.65, 5.9, 0.4, 0.5, "→", 22, GRAY)


# ════════════════════════════════════════════════════════
# Slide 19: Deployment Architecture
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.8, 0.4, 12, 0.7, "部署架构", 36, WHITE, True)
txt(s, 0.8, 1.0, 12, 0.4, "All-in-K8s: 统一架构，在线/离线唯一差异是 LLM 端点 + 镜像仓库", 16, GRAY)

# Deployment scenarios
scenarios = [
    ("开发单机", "1 节点", "single-node RKE2\n本地开发", ACCENT),
    ("小型离线", "1-3 节点", "离线镜像仓库\n+ 本地 LLM", ACCENT2),
    ("中型集群", "3-10 节点", "标准部署\nHA 中间件", GREEN),
    ("云环境", "弹性", "托管 K8s\n+ 云 LLM 端点", ORANGE),
]
for i, (name, scale, desc, clr) in enumerate(scenarios):
    x = 0.8 + i * 3.1
    card(s, x, 1.5, 2.8, 1.8)
    bar(s, x, 1.5, 2.8, 0.06, clr)
    txt(s, x + 0.2, 1.65, 2.4, 0.3, name, 18, clr, True)
    txt(s, x + 0.2, 2.0, 2.4, 0.3, scale, 14, WHITE)
    txt(s, x + 0.2, 2.35, 2.4, 0.7, desc, 14, LIGHT)

# Infrastructure components
txt(s, 0.8, 3.6, 12, 0.4, "基础设施组件", 20, WHITE, True)
infra = [
    ("dts-operator", "Go", "K8s Operator\n组件生命周期\nAppPack 部署", GREEN),
    ("dts-cli", "Go", "CLI 安装/升级\n备份/离线打包\nPack 管理", ACCENT),
    ("Helm Chart", "分层", "platform-base\ndts-core\nextensions", ACCENT2),
    ("Argo CD", "GitOps", "环境级\n持续交付", ORANGE),
]
for i, (name, lang, desc, clr) in enumerate(infra):
    x = 0.8 + i * 3.1
    card(s, x, 4.0, 2.8, 1.6)
    txt(s, x + 0.2, 4.05, 2.0, 0.3, name, 16, clr, True)
    txt(s, x + 2.0, 4.05, 0.7, 0.3, lang, 12, GRAY)
    txt(s, x + 0.2, 4.4, 2.4, 1.0, desc, 14, LIGHT)

# Storage layer
txt(s, 0.8, 5.9, 12, 0.4, "存储层 — 六大引擎各司其职", 18, WHITE, True)
storage = [
    ("PostgreSQL", "主关系存储"),
    ("ClickHouse", "时序分析"),
    ("Neo4j CE", "图关系"),
    ("pgvector", "向量索引"),
    ("MinIO", "对象存储"),
    ("Kafka KRaft", "消息事件"),
]
for i, (name, role) in enumerate(storage):
    x = 0.8 + i * 2.05
    card(s, x, 6.3, 1.85, 0.7)
    txt(s, x + 0.1, 6.33, 1.65, 0.25, name, 13, ACCENT, True, PP_ALIGN.CENTER)
    txt(s, x + 0.1, 6.58, 1.65, 0.3, role, 12, GRAY, False, PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
# Slide 20: Technical Specifications
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.8, 0.4, 12, 0.7, "技术规格", 36, WHITE, True)

# Performance
txt(s, 0.8, 1.1, 6, 0.4, "性能指标", 20, WHITE, True)
perf = [
    ("NL2SQL 响应", "< 3s (P95)"),
    ("简单查询端到端", "< 5s (P95)"),
    ("多步决策任务", "< 120s"),
    ("Skill 调用延迟", "< 500ms (规则) / < 3s (LLM)"),
    ("并发对话数", "100+ (单集群)"),
    ("Ontology 规模", "1000+ ObjectType"),
    ("审计日志吞吐", "10000+ events/s"),
]
for i, (metric, value) in enumerate(perf):
    y = 1.5 + i * 0.45
    card(s, 0.8, y, 5.8, 0.38)
    txt(s, 1.0, y + 0.02, 2.8, 0.3, metric, 14, LIGHT)
    txt(s, 3.8, y + 0.02, 2.5, 0.3, value, 14, GREEN, True)

# Hardware requirements
txt(s, 7.0, 1.1, 6, 0.4, "部署要求", 20, WHITE, True)
hw = [
    ("CPU", "16 核 (最低) / 32 核+ (推荐)"),
    ("内存", "32 GB (最低) / 64 GB+ (推荐)"),
    ("磁盘", "200 GB SSD / 500 GB+ SSD"),
    ("K8s 节点", "1 (开发) / 3+ (生产)"),
    ("LLM", "本地 7B / 云端 GPT-4 级"),
    ("浏览器", "Chrome/Firefox/Edge 100+"),
]
for i, (item, spec) in enumerate(hw):
    y = 1.5 + i * 0.45
    card(s, 7.0, y, 5.8, 0.38)
    txt(s, 7.2, y + 0.02, 1.5, 0.3, item, 14, ACCENT)
    txt(s, 8.7, y + 0.02, 3.8, 0.3, spec, 14, LIGHT)

# Observability
txt(s, 0.8, 4.8, 12, 0.4, "可观测性 — OpenTelemetry 全链路", 20, WHITE, True)
obs = [
    ("Grafana Tempo", "分布式追踪\ntraceId 串联全链路", ACCENT),
    ("Prometheus", "指标监控\n延迟/吞吐/错误率", ACCENT2),
    ("Grafana Loki", "日志聚合检索\n多维度过滤", GREEN),
    ("AI Trace 树", "Agent 推理可视化\nToken 成本统计", ORANGE),
]
for i, (name, desc, clr) in enumerate(obs):
    x = 0.8 + i * 3.1
    card(s, x, 5.2, 2.8, 1.4)
    txt(s, x + 0.2, 5.25, 2.4, 0.3, name, 16, clr, True)
    txt(s, x + 0.2, 5.6, 2.4, 0.8, desc, 14, LIGHT)


# ════════════════════════════════════════════════════════
# Slide 21: 25 Services Overview
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.8, 0.4, 12, 0.7, "25 个微服务全景", 36, WHITE, True)

# Python services
txt(s, 0.8, 1.1, 6, 0.35, "Python AI 核心 (10 个, 全部可降级)", 16, ACCENT, True)
py_services = [
    "dts-agent — Agent 运行时",
    "dts-intent-engine — 意图解析",
    "dts-ontology-engine — 语义推理",
    "dts-data-connector — 数据接入",
    "dts-data-quality — 质量检测",
    "dts-query-ai — NL2SQL",
    "dts-ai-eval — 评测套件",
    "dts-observability — 链路追踪",
    "dts-scheduler — 任务调度",
    "dts-perception — 多模态 (P2)",
]
for i, svc in enumerate(py_services):
    y = 1.5 + i * 0.35
    txt(s, 0.8, y, 5.8, 0.3, f"  {svc}", 12, LIGHT)

# Java services
txt(s, 7.0, 1.1, 6, 0.35, "Java 业务平台 (10 个, 独立可用)", 16, ACCENT2, True)
java_services = [
    "dts-platform — 用户/权限/租户",
    "dts-gateway — API 网关/认证",
    "dts-ontology-store — Ontology CRUD",
    "dts-query-service — SQL 执行",
    "dts-governance — 元数据/血缘",
    "dts-asset — 资产目录",
    "dts-data-security — 动态脱敏",
    "dts-workflow — 审批引擎",
    "dts-audit-log — 审计日志",
    "dts-data-service — 数据 API",
]
for i, svc in enumerate(java_services):
    y = 1.5 + i * 0.35
    txt(s, 7.0, y, 5.8, 0.3, f"  {svc}", 12, LIGHT)

# Go + Frontend
txt(s, 0.8, 5.3, 6, 0.35, "Go 基础设施 (2 个)", 16, GREEN, True)
txt(s, 0.8, 5.7, 5.8, 0.3, "  dts-operator — K8s Operator", 12, LIGHT)
txt(s, 0.8, 6.0, 5.8, 0.3, "  dts-cli — CLI 工具", 12, LIGHT)

txt(s, 7.0, 5.3, 6, 0.35, "Frontend (3 个)", 16, ORANGE, True)
txt(s, 7.0, 5.7, 5.8, 0.3, "  dts-pilot-webapp — 决策门户", 12, LIGHT)
txt(s, 7.0, 6.0, 5.8, 0.3, "  dts-cortex-webapp — 专家工作台", 12, LIGHT)
txt(s, 7.0, 6.3, 5.8, 0.3, "  dts-admin-webapp — 管理控制台", 12, LIGHT)

card(s, 0.8, 6.8, 11.7, 0.5, DARK_BLUE)
txt(s, 1.2, 6.85, 11, 0.35,
    "关键约束: Python 服务是增强层，Java 服务必须能独立运行 — 关掉所有 Python 服务，Java 平台仍然 100% 可用", 14, LIGHT)


# ════════════════════════════════════════════════════════
# Slide 22: Summary
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
bar(s, 1, 1.8, 1.5, 0.06, ACCENT)
txt(s, 1, 2.1, 11, 1.0, "DTS 产品功能说明书", 44, WHITE, True)
txt(s, 1, 3.0, 11, 0.6, "AI Decision Operating System — v3.0", 24, ACCENT)

bullets(s, 1, 3.8, 11, 3.0, [
    "核心公式: Ontology (认知基座) + Intent (意图理解) + Agent (自主执行)",
    "三支柱架构: dts-studio (设计中心) + dts-stack (能力中心) + app-stack (应用中心)",
    "25 个微服务 (Python 10 + Java 10 + Go 2 + Frontend 3)，三语言解耦协作",
    "DAP 五层协议: Ontology → Skill → Intent → Security → Audit，产品核心护城河",
    "57 个 AI 技能: 17 平台 + 6 数据 + 9 AI + 18 行业 + 7 运维",
    "五条铁律: 人工可接管 | 安全不可绕过 | 数据安全基因 | 全操作可追溯 | 能力先于界面",
    "All-in-K8s 统一部署，在线/离线唯一差异是 LLM 端点",
], 18, LIGHT)

txt(s, 1, 6.8, 11, 0.4,
    "把专家脑中的隐性经验变成系统里能自动运转的能力包", 20, ACCENT2, True)


# ── Save ──
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "DTS-产品功能说明书-V1.0.0.pptx")
prs.save(out)
print(f"Saved: {out}")
